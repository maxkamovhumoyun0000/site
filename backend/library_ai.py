"""Teacher library (nested folders) + AI-checked test runtime + weekly review.

Bu modul uchta katta blokni beradi:

1. `/teacher/library/*` — cheksiz chuqurlikdagi papka daraxti. Har bir tugun
   `folder`, `file` yoki `test` bo'lishi mumkin. Tugunlar boshqa o'qituvchilarga
   `view` / `assign` / `edit` huquqi bilan share qilinadi (huquq subtree'ga meros).

2. `/teacher/library/ai/import-screenshot` — kitob sahifasining skrinshotini
   Grok vision bilan o'qib, matn + mashqlarni tayyor test savollariga aylantiradi.
   O'qituvchi tahrirlab saqlaydi.

3. `/student/ai-tests/*` — yangi test turlari runtime'i. Taymer yo'q, har bir
   javob darhol tekshiriladi, xato bo'lsa shu savol tugatilmaguncha keyingisiga
   o'tilmaydi, testdan chiqib ketilsa attempt bekor bo'ladi (boshqatdan).
"""

from __future__ import annotations

import asyncio
import json
import logging
import mimetypes
import os
import random
import re
from datetime import datetime, timedelta
from pathlib import Path
from typing import Any
from urllib.parse import urlparse

from fastapi import APIRouter, File, Header, HTTPException, Query, UploadFile
from pydantic import BaseModel, Field

import db as dbm

router = APIRouter()
logger = logging.getLogger(__name__)


# ═══════════════════════════════════════════════════════════════════════════
# Test turlari registri — sayt va ikkala ilova shu ro'yxatdan foydalanadi
# ═══════════════════════════════════════════════════════════════════════════

#: `check` — javob qanday tekshiriladi: "auto" (server solishtiradi) yoki "ai".
#: `input` — student nima yuboradi: "text" | "audio" | "audio_or_text" | "choice" | "order" | "pairs"
#: `needs_audio_asset` — o'qituvchi mashqqa audio fayl yuklashi shart.
AI_TEST_TYPES: dict[str, dict[str, Any]] = {
    # ── Foydalanuvchi so'ragan asosiy turlar ────────────────────────────────
    "speak_sentence": {
        "label_uz": "So'z bilan gap tuzib gapirish",
        "label_ru": "Составить предложение и произнести",
        "label_en": "Speak a sentence",
        "check": "ai", "input": "audio", "needs_audio_asset": False,
        "retry_until_correct": True,
    },
    "write_sentence": {
        "label_uz": "So'z bilan gap tuzib yozish",
        "label_ru": "Составить предложение письменно",
        "label_en": "Write a sentence",
        "check": "ai", "input": "text", "needs_audio_asset": False,
        "retry_until_correct": True,
    },
    "guided_writing": {
        "label_uz": "Mavzu bo'yicha yozma mashq",
        "label_ru": "Письменное задание по теме",
        "label_en": "Guided writing",
        "check": "ai", "input": "text", "needs_audio_asset": False,
        "retry_until_correct": True,
    },
    # ── Kitob mashqlarida uchraydigan qo'shimcha turlar (AI tekshiradi) ─────
    "translation": {
        "label_uz": "Tarjima qilish",
        "label_ru": "Перевод",
        "label_en": "Translation",
        "check": "ai", "input": "text", "needs_audio_asset": False,
        "retry_until_correct": True,
    },
    "reading_open": {
        "label_uz": "Matn bo'yicha ochiq savol",
        "label_ru": "Открытый вопрос по тексту",
        "label_en": "Reading comprehension (open)",
        "check": "ai", "input": "text", "needs_audio_asset": False,
        "retry_until_correct": True,
    },
    "read_aloud": {
        "label_uz": "Matnni ovoz chiqarib o'qish",
        "label_ru": "Чтение вслух",
        "label_en": "Read aloud",
        "check": "ai", "input": "audio", "needs_audio_asset": False,
        "retry_until_correct": True,
    },
    "paraphrase": {
        "label_uz": "Gapni boshqacha aytish",
        "label_ru": "Перефразировать",
        "label_en": "Paraphrase",
        "check": "ai", "input": "text", "needs_audio_asset": False,
        "retry_until_correct": True,
    },
    "dialogue_completion": {
        "label_uz": "Dialogni to'ldirish",
        "label_ru": "Дополнить диалог",
        "label_en": "Dialogue completion",
        "check": "ai", "input": "audio_or_text", "needs_audio_asset": False,
        "retry_until_correct": True,
    },
    "picture_description": {
        "label_uz": "Rasmni tasvirlash",
        "label_ru": "Описать картинку",
        "label_en": "Picture description",
        "check": "ai", "input": "audio_or_text", "needs_audio_asset": False,
        "retry_until_correct": True,
    },
    # ── Avtomatik tekshiriladigan turlar (AI chaqirilmaydi, tez ishlaydi) ───
    "listening": {
        "label_uz": "Tinglab tushunish",
        "label_ru": "Аудирование",
        "label_en": "Listening",
        "check": "auto", "input": "choice", "needs_audio_asset": True,
        "retry_until_correct": False,
    },
    "dictation": {
        "label_uz": "Diktant (tinglab yozish)",
        "label_ru": "Диктант",
        "label_en": "Dictation",
        "check": "auto", "input": "text", "needs_audio_asset": True,
        "retry_until_correct": True,
    },
    "spelling": {
        "label_uz": "So'zni to'g'ri yozish",
        "label_ru": "Правописание",
        "label_en": "Spelling",
        "check": "auto", "input": "text", "needs_audio_asset": False,
        "retry_until_correct": True,
    },
    "matching": {
        "label_uz": "Juftlab moslashtirish",
        "label_ru": "Соотнести пары",
        "label_en": "Matching",
        "check": "auto", "input": "pairs", "needs_audio_asset": False,
        "retry_until_correct": False,
    },
    "scrambled_sentence": {
        "label_uz": "So'zlarni tartibga solish",
        "label_ru": "Составить предложение из слов",
        "label_en": "Scrambled sentence",
        "check": "auto", "input": "order", "needs_audio_asset": False,
        "retry_until_correct": True,
    },
    "gap_fill": {
        "label_uz": "Bo'sh joyni to'ldirish",
        "label_ru": "Заполнить пропуск",
        "label_en": "Gap fill",
        "check": "auto", "input": "text", "needs_audio_asset": False,
        "retry_until_correct": True,
    },
    # ── Yaxlit matn + bir nechta bo'sh joy + so'zlar banki (Duolingo/kitob uslubi) ──
    "passage_cloze": {
        "label_uz": "Matnni to'ldirish (so'zlar banki bilan)",
        "label_ru": "Заполнить текст (с банком слов)",
        "label_en": "Passage cloze (with word bank)",
        "check": "auto", "input": "cloze", "needs_audio_asset": False,
        "retry_until_correct": True,
    },
    # ── O'qish matni + bir nechta turli savol (true/false/NG, sinonim, gap, kim aytdi) ──
    "reading_set": {
        "label_uz": "Matn va savollar",
        "label_ru": "Текст и вопросы",
        "label_en": "Reading set",
        "check": "auto", "input": "reading_set", "needs_audio_asset": False,
        "retry_until_correct": True,
    },
    # ── Polimorf: bitta so'z. Studentga tushganda random test turiga aylanadi ──
    "word_practice": {
        "label_uz": "So'z mashqi (random tur)",
        "label_ru": "Практика слова (случайный тип)",
        "label_en": "Word practice (random type)",
        "check": "ai", "input": "text", "needs_audio_asset": False,
        "retry_until_correct": True,
        "polymorphic": True,
    },
}

AI_CHECKED_KINDS = {k for k, v in AI_TEST_TYPES.items() if v["check"] == "ai"}
AUTO_CHECKED_KINDS = {k for k, v in AI_TEST_TYPES.items() if v["check"] == "auto"}
AUDIO_ASSET_KINDS = {k for k, v in AI_TEST_TYPES.items() if v["needs_audio_asset"]}
POLYMORPHIC_KINDS = {k for k, v in AI_TEST_TYPES.items() if v.get("polymorphic")}

#: word_practice studentga tushganda shu turlardan biriga random aylanadi.
WORD_PRACTICE_VARIANTS = ["speak_sentence", "write_sentence", "spelling", "translation"]

#: AI turli nom bilan qaytarishi mumkin — ularni bizning kanonik turlarga moslaymiz.
_KIND_SYNONYMS: dict[str, str] = {
    "multiple_choice": "listening",
    "mcq": "listening",
    "choice": "listening",
    "true_false": "listening",
    "fill_in_the_blank": "gap_fill",
    "fill_blank": "gap_fill",
    "fill_gap": "gap_fill",
    "cloze": "gap_fill",
    "gapfill": "gap_fill",
    "gap": "gap_fill",
    "match": "matching",
    "match_pairs": "matching",
    "matching_pairs": "matching",
    "reorder": "scrambled_sentence",
    "order_words": "scrambled_sentence",
    "unscramble": "scrambled_sentence",
    "word_order": "scrambled_sentence",
    "sentence_building": "scrambled_sentence",
    "cloze": "passage_cloze",
    "cloze_passage": "passage_cloze",
    "passage_gap_fill": "passage_cloze",
    "text_completion": "passage_cloze",
    "complete_the_text": "passage_cloze",
    "complete_the_passage": "passage_cloze",
    "fill_passage": "passage_cloze",
    "gapped_text": "passage_cloze",
    "reading_set": "reading_set",
    "reading_passage": "reading_set",
    "reading_questions": "reading_set",
    "text_with_questions": "reading_set",
    "passage_questions": "reading_set",
    "vocabulary": "word_practice",
    "vocab": "word_practice",
    "word": "word_practice",
    "spelling_word": "spelling",
    "dictation_audio": "dictation",
    "listening_comprehension": "listening",
    "reading": "reading_open",
    "reading_comprehension": "reading_open",
    "comprehension": "reading_open",
    "open_question": "reading_open",
    "short_answer": "reading_open",
    "speaking": "speak_sentence",
    "speak": "speak_sentence",
    "writing": "guided_writing",
    "essay": "guided_writing",
    "translate": "translation",
    "picture": "picture_description",
    "describe_picture": "picture_description",
    "dialogue": "dialogue_completion",
    "conversation": "dialogue_completion",
}

MAX_RETRIES_PER_QUESTION = 6


@router.get("/test-types")
async def list_test_types(x_language: str | None = Header(default=None, alias="X-Language")):
    """Sayt va ilovalar test turlari ro'yxatini shu yerdan oladi (bir manba)."""
    lang = str(x_language or "uz").strip().lower()[:2]
    key = f"label_{lang}" if f"label_{lang}" in next(iter(AI_TEST_TYPES.values())) else "label_uz"
    return {
        "items": [
            {
                "kind": kind,
                "label": meta.get(key) or meta.get("label_uz"),
                "check": meta["check"],
                "input": meta["input"],
                "needs_audio_asset": meta["needs_audio_asset"],
                "retry_until_correct": meta["retry_until_correct"],
            }
            for kind, meta in AI_TEST_TYPES.items()
        ]
    }


# ═══════════════════════════════════════════════════════════════════════════
# Auth yordamchilari (backend.main bilan tsiklik importdan qochish uchun lazy)
# ═══════════════════════════════════════════════════════════════════════════

def _auth(authorization: str | None, roles: set[str]) -> dict:
    from backend.main import _require_role, _user_row_from_bearer

    user = _user_row_from_bearer(authorization)
    _require_role(user, roles)
    return user


TEACHER_ROLES = {"teacher", "support", "admin", "superadmin"}
STUDENT_ROLES = {"student"}


def _settings() -> dict[str, float]:
    from backend.main import _get_runtime_settings

    try:
        return _get_runtime_settings() or {}
    except Exception:
        logger.exception("dpoint settings read failed")
        return {}


def _setting(key: str, default: float) -> float:
    try:
        return float(_settings().get(key, default))
    except Exception:
        return float(default)


def _safe(fn, fallback=None):
    try:
        return fn()
    except HTTPException:
        raise
    except Exception:
        logger.exception("library_ai safe call failed")
        return fallback


# ═══════════════════════════════════════════════════════════════════════════
# 1. KUTUBXONA DARAXTI
# ═══════════════════════════════════════════════════════════════════════════

class LibraryNodeCreate(BaseModel):
    title: str = Field(min_length=1, max_length=300)
    kind: str = "folder"
    parent_id: int | None = None
    description: str | None = None
    subject: str | None = None
    level: str | None = None
    file_url: str | None = None
    payload: dict | None = None
    is_public: bool = False


class LibraryNodeUpdate(BaseModel):
    title: str | None = None
    description: str | None = None
    subject: str | None = None
    level: str | None = None
    file_url: str | None = None
    payload: dict | None = None
    is_public: bool | None = None
    parent_id: int | None = None
    sort_order: int | None = None


class LibraryShareRequest(BaseModel):
    teacher_id: int
    permission: str = "view"


class LibraryAssignRequest(BaseModel):
    student_ids: list[int] | None = None
    group_id: int | None = None
    title: str | None = None
    description: str | None = None
    due_at: str | None = None


def _node_or_404(node_id: int) -> dict:
    node = _safe(lambda: dbm.get_library_node(int(node_id)))
    if not node:
        raise HTTPException(status_code=404, detail="Kutubxona elementi topilmadi")
    return node


def _require_node_permission(user: dict, node: dict, need: str) -> str:
    """need: 'view' | 'assign' | 'edit'. Admin hamma narsani ko'radi/tahrirlaydi."""
    from backend.main import _role_from_login_type

    role = _role_from_login_type(int(user.get("login_type") or 1), str(user.get("login_id") or ""))
    if role in {"admin", "superadmin"}:
        return "owner"
    perm = _safe(lambda: dbm.library_permission_for(None, int(user.get("id") or 0), node))
    if perm is None and node.get("is_public") and need == "view":
        perm = "view"
    if perm is None:
        raise HTTPException(status_code=403, detail="Bu elementga ruxsatingiz yo'q")
    rank = {"view": 1, "assign": 2, "edit": 3, "owner": 4}
    if rank.get(perm, 0) < rank.get(need, 99):
        raise HTTPException(status_code=403, detail="Bu amal uchun huquq yetarli emas")
    return perm


@router.get("/teacher/library")
async def library_list(authorization: str | None = Header(default=None)):
    """O'ziniki + public + share qilingan barcha tugunlar (flat, frontend daraxt qiladi)."""
    user = _auth(authorization, TEACHER_ROLES)
    data = _safe(lambda: dbm.list_library_nodes(int(user.get("id") or 0)), {"nodes": [], "shares": []}) or {}
    nodes = []
    for node in data.get("nodes") or []:
        item = dict(node)
        item["payload"] = _json_obj(item.pop("payload_json", None))
        nodes.append(item)
    return {
        "nodes": nodes,
        "shares": data.get("shares") or [],
        "my_id": int(user.get("id") or 0),
        "total": len(nodes),
    }


@router.get("/teacher/library/{node_id}")
async def library_get(node_id: int, authorization: str | None = Header(default=None)):
    user = _auth(authorization, TEACHER_ROLES)
    node = _node_or_404(node_id)
    perm = _require_node_permission(user, node, "view")
    out = dict(node)
    out["payload"] = _json_obj(out.pop("payload_json", None))
    out["permission"] = perm
    return {"item": out}


@router.post("/teacher/library")
async def library_create(payload: LibraryNodeCreate, authorization: str | None = Header(default=None)):
    user = _auth(authorization, TEACHER_ROLES)
    kind = str(payload.kind or "folder").strip().lower()
    if kind not in dbm.LIBRARY_KINDS:
        raise HTTPException(status_code=400, detail="Yaroqsiz element turi")
    if payload.parent_id:
        parent = _node_or_404(int(payload.parent_id))
        if str(parent.get("kind") or "") != "folder":
            raise HTTPException(status_code=400, detail="Faqat papka ichiga element qo'shiladi")
        _require_node_permission(user, parent, "edit")
    if kind == "test":
        questions = _normalize_questions((payload.payload or {}).get("questions"))
        if not questions:
            raise HTTPException(status_code=400, detail="Test uchun kamida bitta savol kerak")
        payload.payload = {**(payload.payload or {}), "questions": questions}
    try:
        node = dbm.create_library_node(
            int(user.get("id") or 0),
            payload.title,
            kind,
            parent_id=payload.parent_id,
            description=payload.description,
            subject=payload.subject,
            level=payload.level,
            file_url=payload.file_url,
            payload=payload.payload,
            is_public=bool(payload.is_public),
        )
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc))
    out = dict(node)
    out["payload"] = _json_obj(out.pop("payload_json", None))
    return {"message": "Element yaratildi", "item": out}


@router.patch("/teacher/library/{node_id}")
async def library_update(node_id: int, payload: LibraryNodeUpdate, authorization: str | None = Header(default=None)):
    user = _auth(authorization, TEACHER_ROLES)
    node = _node_or_404(node_id)
    _require_node_permission(user, node, "edit")
    new_payload = payload.payload
    if new_payload is not None and str(node.get("kind") or "") == "test":
        questions = _normalize_questions(new_payload.get("questions"))
        if not questions:
            raise HTTPException(status_code=400, detail="Test uchun kamida bitta savol kerak")
        new_payload = {**new_payload, "questions": questions}
    if payload.parent_id is not None and int(payload.parent_id or 0) > 0:
        target = _node_or_404(int(payload.parent_id))
        if str(target.get("kind") or "") != "folder":
            raise HTTPException(status_code=400, detail="Faqat papka ichiga ko'chiriladi")
        _require_node_permission(user, target, "edit")
    try:
        updated = dbm.update_library_node(
            int(node_id),
            title=payload.title,
            description=payload.description,
            subject=payload.subject,
            level=payload.level,
            file_url=payload.file_url,
            payload=new_payload,
            is_public=payload.is_public,
            parent_id=payload.parent_id,
            sort_order=payload.sort_order,
        )
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc))
    if not updated:
        raise HTTPException(status_code=404, detail="Element topilmadi")
    out = dict(updated)
    out["payload"] = _json_obj(out.pop("payload_json", None))
    return {"message": "Saqlandi", "item": out}


@router.delete("/teacher/library/{node_id}")
async def library_delete(node_id: int, authorization: str | None = Header(default=None)):
    user = _auth(authorization, TEACHER_ROLES)
    node = _node_or_404(node_id)
    from backend.main import _role_from_login_type

    role = _role_from_login_type(int(user.get("login_type") or 1), str(user.get("login_id") or ""))
    if int(node.get("owner_id") or 0) != int(user.get("id") or 0) and role not in {"admin", "superadmin"}:
        raise HTTPException(status_code=403, detail="Faqat egasi o'chira oladi")
    ok = _safe(lambda: dbm.delete_library_node(int(node_id)), False)
    if not ok:
        raise HTTPException(status_code=404, detail="Element topilmadi")
    return {"message": "O'chirildi"}


@router.get("/teacher/library/{node_id}/shares")
async def library_shares(node_id: int, authorization: str | None = Header(default=None)):
    user = _auth(authorization, TEACHER_ROLES)
    node = _node_or_404(node_id)
    _require_node_permission(user, node, "view")
    data = _safe(lambda: dbm.list_library_share_targets(int(node_id)), {"shares": [], "teachers": []}) or {}
    return data


@router.post("/teacher/library/{node_id}/share")
async def library_share(node_id: int, payload: LibraryShareRequest, authorization: str | None = Header(default=None)):
    user = _auth(authorization, TEACHER_ROLES)
    node = _node_or_404(node_id)
    from backend.main import _role_from_login_type

    role = _role_from_login_type(int(user.get("login_type") or 1), str(user.get("login_id") or ""))
    if int(node.get("owner_id") or 0) != int(user.get("id") or 0) and role not in {"admin", "superadmin"}:
        raise HTTPException(status_code=403, detail="Faqat egasi share qila oladi")
    if int(payload.teacher_id or 0) == int(node.get("owner_id") or 0):
        raise HTTPException(status_code=400, detail="Egasiga share qilish shart emas")
    try:
        result = dbm.share_library_node(
            int(node_id), int(payload.teacher_id), payload.permission, shared_by=int(user.get("id") or 0)
        )
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc))
    return {"message": "Ulashildi", "share": result}


@router.delete("/teacher/library/{node_id}/share/{teacher_id}")
async def library_unshare(node_id: int, teacher_id: int, authorization: str | None = Header(default=None)):
    user = _auth(authorization, TEACHER_ROLES)
    node = _node_or_404(node_id)
    from backend.main import _role_from_login_type

    role = _role_from_login_type(int(user.get("login_type") or 1), str(user.get("login_id") or ""))
    if int(node.get("owner_id") or 0) != int(user.get("id") or 0) and role not in {"admin", "superadmin"}:
        raise HTTPException(status_code=403, detail="Faqat egasi bekor qila oladi")
    dbm.unshare_library_node(int(node_id), int(teacher_id))
    return {"message": "Ulashish bekor qilindi"}


@router.post("/teacher/library/upload")
async def library_upload(
    file: UploadFile = File(...),
    authorization: str | None = Header(default=None),
):
    """Kutubxona uchun umumiy fayl yuklash (PDF, rasm, listening audio)."""
    user = _auth(authorization, TEACHER_ROLES)
    from backend.main import HOMEWORK_UPLOAD_DIR, _upload_web_file

    filename, _ = await _upload_web_file(
        file, HOMEWORK_UPLOAD_DIR, f"library_{int(user.get('id') or 0)}", max_size_mb=25
    )
    return {"url": f"/homework/files/{filename}", "filename": filename}


@router.post("/teacher/library/{node_id}/assign")
async def library_assign(node_id: int, payload: LibraryAssignRequest, authorization: str | None = Header(default=None)):
    """Kutubxonadagi test/faylni studentlarga homework qilib beradi."""
    user = _auth(authorization, TEACHER_ROLES)
    node = _node_or_404(node_id)
    _require_node_permission(user, node, "assign")
    kind = str(node.get("kind") or "")
    if kind == "folder":
        raise HTTPException(status_code=400, detail="Papkani homework qilib berib bo'lmaydi")
    node_payload = _json_obj(node.get("payload_json"))
    questions = _normalize_questions(node_payload.get("questions")) if kind == "test" else []
    if kind == "test" and not questions:
        raise HTTPException(status_code=400, detail="Bu testda savol yo'q")

    student_ids = [int(s) for s in (payload.student_ids or []) if int(s or 0) > 0]
    group_id = int(payload.group_id or 0) or None
    if not student_ids and not group_id:
        raise HTTPException(status_code=400, detail="Student yoki guruh tanlanmadi")

    title = str(payload.title or node.get("title") or "Kutubxona vazifasi").strip()
    description = str(payload.description or node.get("description") or "").strip() or None
    created: list[dict] = []
    targets: list[int | None] = student_ids or [None]
    for student_id in targets:
        homework = _safe(
            lambda sid=student_id: dbm.create_homework(
                int(user.get("id") or 0),
                sid,
                title,
                description=description,
                due_at=payload.due_at,
                group_id=group_id if sid is None else None,
                homework_kind="test" if kind == "test" else "list",
            )
        )
        if not homework:
            continue
        homework_id = int(homework.get("id") or 0)
        if kind == "test" and homework_id > 0:
            _safe(
                lambda hid=homework_id: dbm.save_content_test(
                    "homework",
                    hid,
                    json.dumps(questions, ensure_ascii=False),
                    int(user.get("id") or 0),
                    title=title,
                    raw_questions=True,
                )
            )
        created.append(homework)
    if not created:
        raise HTTPException(status_code=500, detail="Vazifa yaratilmadi")

    # O'quvchilarga bildirishnoma (web + push) — oddiy homework kabi.
    try:
        from backend.main import _notify_homework_created, _homework_audience_student_ids

        for hw in created:
            audience = _homework_audience_student_ids(hw)
            if audience:
                await _notify_homework_created(hw, user, audience)
    except Exception:
        logger.exception("library assign notification failed node_id=%s", node_id)

    return {"message": f"{len(created)} ta vazifa berildi", "items": created}


def _json_obj(raw: Any) -> dict:
    if isinstance(raw, dict):
        return raw
    if not raw:
        return {}
    try:
        parsed = json.loads(str(raw))
        return parsed if isinstance(parsed, dict) else {}
    except Exception:
        return {}


def _normalize_questions(raw: Any) -> list[dict]:
    """Yangi test turlari uchun savollarni normallashtiradi va validatsiya qiladi."""
    if isinstance(raw, str):
        try:
            raw = json.loads(raw)
        except Exception:
            raw = []
    if not isinstance(raw, list):
        return []
    out: list[dict] = []
    for item in raw:
        if not isinstance(item, dict):
            continue
        kind = str(item.get("kind") or item.get("question_type") or item.get("type") or "").strip().lower()
        kind = _KIND_SYNONYMS.get(kind, kind)
        if kind not in AI_TEST_TYPES:
            # Turi noaniq bo'lsa — mazmuniga qarab eng mos turini tanlaymiz
            # (hech qachon jimgina tashlab yubormaymiz).
            if item.get("pairs"):
                kind = "matching"
            elif item.get("options"):
                kind = "listening"
            elif item.get("tokens"):
                kind = "scrambled_sentence"
            elif item.get("word") and not item.get("answer"):
                kind = "word_practice"
            elif str(item.get("prompt") or item.get("question") or "").find("___") >= 0:
                kind = "gap_fill"
            elif item.get("passage"):
                kind = "reading_open"
            elif item.get("answer"):
                kind = "gap_fill"
            else:
                kind = "write_sentence"
        meta = AI_TEST_TYPES[kind]
        question: dict[str, Any] = {
            "kind": kind,
            "prompt": str(item.get("prompt") or item.get("question") or "").strip(),
            "instruction": str(item.get("instruction") or "").strip() or None,
            "word": str(item.get("word") or "").strip() or None,
            "passage": str(item.get("passage") or "").strip() or None,
            "image_url": str(item.get("image_url") or "").strip() or None,
            "audio_url": str(item.get("audio_url") or "").strip() or None,
            "level": str(item.get("level") or "").strip() or None,
            "topic": str(item.get("topic") or "").strip() or None,
            # Bir mashq blokiga tegishli savollarni birlashtirish uchun (AI beradi).
            "group": str(item.get("group") or item.get("exercise_id") or item.get("group_id") or "").strip() or None,
        }
        if not question["prompt"] and not question["word"] and kind != "passage_cloze":
            continue
        if meta["needs_audio_asset"] and not question["audio_url"]:
            # Audio yuklanmagan listening/diktant savoli studentga ko'rsatilmaydi
            question["needs_audio_upload"] = True
        if meta["check"] == "auto":
            if kind == "passage_cloze":
                normalized_cloze = _normalize_passage_cloze(item)
                if not normalized_cloze:
                    continue
                question.update(normalized_cloze)
            elif kind == "reading_set":
                normalized_set = _normalize_reading_set(item)
                if not normalized_set:
                    continue
                question.update(normalized_set)
            elif kind == "matching":
                pairs = item.get("pairs")
                clean_pairs = []
                if isinstance(pairs, list):
                    for pair in pairs:
                        if isinstance(pair, dict):
                            left = str(pair.get("left") or "").strip()
                            right = str(pair.get("right") or "").strip()
                            if left and right:
                                clean_pairs.append({"left": left, "right": right})
                if len(clean_pairs) < 2:
                    continue
                question["pairs"] = clean_pairs
            elif kind in {"scrambled_sentence"}:
                answer = str(item.get("answer") or "").strip()
                if not answer:
                    continue
                question["answer"] = answer
                tokens = item.get("tokens")
                question["tokens"] = (
                    [str(t).strip() for t in tokens if str(t).strip()]
                    if isinstance(tokens, list) and tokens
                    else answer.replace(".", "").split()
                )
                # AI qo'shimcha (chalg'ituvchi) so'zlarni ham beradi — ular pulga aralashtiriladi.
                distractors = item.get("distractors") or item.get("extra_words")
                question["distractors"] = (
                    [str(d).strip() for d in distractors if str(d).strip()]
                    if isinstance(distractors, list) else []
                )
            elif kind == "listening":
                options = [str(o).strip() for o in (item.get("options") or []) if str(o).strip()]
                if len(options) < 2:
                    continue
                question["options"] = options[:4]
                try:
                    question["correct_index"] = max(0, min(len(question["options"]) - 1, int(item.get("correct_index") or item.get("correct_option_index") or 0)))
                except Exception:
                    question["correct_index"] = 0
            else:  # dictation, spelling, gap_fill
                answer = str(item.get("answer") or "").strip()
                if not answer:
                    continue
                question["answer"] = answer
                accepted = item.get("accepted_answers")
                question["accepted_answers"] = (
                    [str(a).strip() for a in accepted if str(a).strip()] if isinstance(accepted, list) else []
                )
        else:
            question["reference_answer"] = str(item.get("reference_answer") or item.get("answer") or "").strip() or None
            question["target_level"] = str(item.get("target_level") or item.get("level") or "").strip() or None
            if kind == "word_practice":
                # So'z mashqi: studentga tushganda random turga aylanadi.
                word = question.get("word") or question.get("prompt")
                if not word:
                    continue
                question["word"] = word
                # Har bir so'z uchun ikkala tarjima (uz + ru) bir testda saqlanadi.
                question["translation_uz"] = str(item.get("translation_uz") or "").strip() or None
                question["translation_ru"] = str(item.get("translation_ru") or "").strip() or None
                # Backward-compat.
                question["translation"] = (
                    str(item.get("translation") or "").strip()
                    or question["translation_uz"] or question["translation_ru"]
                )
                question["meaning"] = str(item.get("meaning") or "").strip() or question["translation_ru"]
        out.append(question)
    return out


# Matn ichidagi bo'sh joy belgilari: ___ (3+ tag chiziq), [[1]], {{1}}, (1) ______
_CLOZE_GAP_RE = re.compile(r"\{\{\s*\d+\s*\}\}|\[\[\s*\d+\s*\]\]|_{2,}|\.{4,}")


def _normalize_passage_cloze(item: dict) -> dict | None:
    """Yaxlit matn + bir nechta bo'sh joy + so'zlar bankini normallashtiradi.

    Bir necha kirish formatini qabul qiladi:
      • passage matnida ___ / {{1}} / [[1]] belgilari + `answers`/`blanks` ro'yxati
      • `sentences`: har biri bitta gap+javob (ular yagona matnga birlashtiriladi)
    """
    passage = str(item.get("passage") or item.get("text") or item.get("prompt") or "").strip()
    # Javoblar: answers / blanks / gaps
    raw_answers = item.get("answers") or item.get("blanks") or item.get("gaps") or []
    answers: list[dict] = []
    if isinstance(raw_answers, list):
        for a in raw_answers:
            if isinstance(a, dict):
                ans = str(a.get("answer") or a.get("value") or "").strip()
                acc = a.get("accepted_answers") or a.get("accepted") or []
                accepted = [str(x).strip() for x in acc if str(x).strip()] if isinstance(acc, list) else []
            else:
                ans = str(a or "").strip()
                accepted = []
            if ans:
                answers.append({"answer": ans, "accepted_answers": accepted})

    # `sentences` formatidan yagona matn quramiz.
    if not passage and isinstance(item.get("sentences"), list):
        lines = []
        for s in item["sentences"]:
            if not isinstance(s, dict):
                continue
            text = str(s.get("text") or s.get("sentence") or "").strip()
            ans = str(s.get("answer") or "").strip()
            if not text or not ans:
                continue
            if not _CLOZE_GAP_RE.search(text):
                text = f"{text} ___"
            lines.append(text)
            answers.append({"answer": ans, "accepted_answers": []})
        passage = "\n".join(lines)

    if not passage or not answers:
        return None

    # Matndagi bo'sh joylarni yagona `___` belgisiga keltiramiz.
    normalized_passage = _CLOZE_GAP_RE.sub("___", passage)
    gap_count = normalized_passage.count("___")
    if gap_count == 0:
        # Belgilar topilmadi — javoblar soniga qarab oxiriga qo'shib qo'yamiz.
        normalized_passage = normalized_passage.rstrip() + " " + " ".join(["___"] * len(answers))
        gap_count = len(answers)
    # Bo'sh joy soni bilan javob soni mos kelishi kerak.
    if gap_count != len(answers):
        answers = answers[:gap_count] if len(answers) > gap_count else answers + [
            {"answer": "", "accepted_answers": []} for _ in range(gap_count - len(answers))
        ]
    answers = [a for a in answers if a.get("answer")]
    if not answers:
        return None

    # Alohida gapli mashqlar (11.1 kabi) raqamlanadi: "1. ... 2. ...".
    # Uzluksiz o'qish matni (11.3 kabi) raqamlanmaydi — oqma matn bo'lib qoladi.
    normalized_passage = _number_cloze_sentences(normalized_passage, len(answers))

    # So'zlar banki: AI bergan qutini (box) aynan olamiz — chalg'ituvchi so'z
    # QO'SHMAYMIZ. Faqat qutida hech narsa bo'lmasa, javoblarning o'zini beramiz.
    bank_raw = item.get("word_bank") or item.get("bank") or item.get("options") or []
    bank = [str(w).strip() for w in bank_raw if str(w).strip()] if isinstance(bank_raw, list) else []
    if not bank:
        bank = [a["answer"] for a in answers]
    seen: set[str] = set()
    bank = [w for w in bank if not (w.lower() in seen or seen.add(w.lower()))]
    import random as _r
    _r.shuffle(bank)

    return {
        # Muharrir uchun: passage (___ bilan) + answers massivi.
        "passage": normalized_passage,
        "answers": answers,
        # Student runner uchun.
        "passage_template": normalized_passage,
        "blanks": answers,
        "word_bank": bank,
    }


#: reading_set ichidagi savol turlari.
READING_SUBTYPES = {"true_false_ng", "choice", "gap", "short", "synonym", "who_said"}


def _normalize_reading_set(item: dict) -> dict | None:
    """Matn + bir nechta turli savol (bitta kartada) mashqini normallashtiradi."""
    passage = str(item.get("passage") or item.get("text") or item.get("reading_text") or "").strip()
    raw_qs = item.get("questions") or item.get("items") or item.get("sub_questions") or []
    if not passage or not isinstance(raw_qs, list):
        return None
    subs: list[dict] = []
    for raw in raw_qs:
        if not isinstance(raw, dict):
            continue
        prompt = str(raw.get("prompt") or raw.get("question") or "").strip()
        answer = str(raw.get("answer") or "").strip()
        if not prompt or not answer:
            continue
        stype = str(raw.get("type") or raw.get("subtype") or "").strip().lower()
        options = [str(o).strip() for o in (raw.get("options") or []) if str(o).strip()]
        if stype in {"true_false", "tfng", "true_false_not_given"}:
            stype = "true_false_ng"
        if stype not in READING_SUBTYPES:
            stype = "choice" if options else "short"
        if stype == "true_false_ng" and not options:
            options = ["True", "False", "Not given"]
        accepted = raw.get("accepted_answers") or []
        subs.append({
            "type": stype,
            "prompt": prompt,
            "options": options,
            "answer": answer,
            "accepted_answers": [str(a).strip() for a in accepted if str(a).strip()],
        })
    if not subs:
        return None
    return {"passage": passage, "sub_questions": subs}


def _number_cloze_sentences(passage: str, gap_count: int) -> str:
    """Alohida gapli cloze mashqlarini raqamlab, har birini yangi qatorga chiqaradi.

    11.1 kabi mashqlar (har gapda bitta bo'sh joy) -> "1. ... \n2. ...".
    11.3 kabi uzluksiz o'qish matni (gaplar soni bo'sh joylar sonidan ancha ko'p
    yoki matn bitta hikoya) -> tegilmaydi.
    """
    text = str(passage or "")
    if gap_count < 2:
        return text
    # Allaqachon raqamlangan bo'lsa (1. / 1) ) tegilmaydi.
    if re.search(r"(^|\n)\s*\d+\s*[.)]\s", text):
        return text

    lines = [ln.strip() for ln in text.split("\n") if ln.strip()]
    if len(lines) > 1:
        # Har qatorda bittadan gap — shunchaki raqamlaymiz.
        if all("___" in ln for ln in lines):
            return "\n".join(f"{i + 1}. {ln}" for i, ln in enumerate(lines))
        return text

    # Bitta uzun qator: gap chegaralari bo'yicha bo'lamiz.
    single = lines[0] if lines else ""
    parts = [p.strip() for p in re.split(r"(?<=[.!?])\s+(?=[A-ZА-ЯЁ“\"'])", single) if p.strip()]
    gapped = [p for p in parts if "___" in p]
    # Har gapda kamida bitta bo'sh joy bo'lsa va gaplar soni bo'sh joylar soniga
    # yaqin bo'lsa — bu alohida gaplar to'plami (mashq), matn emas.
    if len(gapped) >= 2 and len(gapped) == len(parts) and len(parts) >= max(2, gap_count // 2):
        return "\n".join(f"{i + 1}. {p}" for i, p in enumerate(parts))
    return text


def _cloze_from_gap_fill_run(run: list[dict]) -> dict:
    """Bir nechta gap_fill savolni bitta yaxlit passage_cloze ga birlashtiradi."""
    import random as _r

    lines: list[str] = []
    answers: list[dict] = []
    bank: list[str] = []
    instruction = None
    idx = 0
    for q in run:
        prompt = str(q.get("prompt") or "").strip()
        ans = str(q.get("answer") or "").strip()
        if not ans:
            continue
        if "___" not in prompt:
            prompt = (prompt + " ___").strip() if prompt else "___"
        idx += 1
        # Har mashqni alohida raqamli qatorda ko'rsatamiz (chalkash bo'lmasin).
        lines.append(f"{idx}. {prompt}")
        acc = q.get("accepted_answers") or []
        answers.append({"answer": ans, "accepted_answers": [str(a).strip() for a in acc if str(a).strip()]})
        if not instruction and q.get("instruction"):
            instruction = q.get("instruction")
        bank.append(ans)
    passage = "\n".join(lines)
    seen: set[str] = set()
    bank = [w for w in bank if not (w.lower() in seen or seen.add(w.lower()))]
    _r.shuffle(bank)
    return {
        "kind": "passage_cloze",
        "instruction": instruction,
        "passage": passage,
        "answers": answers,
        "passage_template": passage,
        "blanks": answers,
        "word_bank": bank,
    }


def _propagate_group_instructions(questions: list[dict]) -> list[dict]:
    """Bir mashq blokidagi (group) barcha savollarga o'sha mashq shartini beradi.

    11.4 kabi mashqlar alohida savollar bo'lib chiqadi — har birining tepasida
    umumiy shart ("Write sentences about the past") ko'rinishi kerak."""
    by_group: dict[str, str] = {}
    for q in questions or []:
        gid = str(q.get("group") or "").strip()
        ins = str(q.get("instruction") or "").strip()
        if gid and ins and gid not in by_group:
            by_group[gid] = ins
    if not by_group:
        return questions
    for q in questions:
        gid = str(q.get("group") or "").strip()
        if gid and not str(q.get("instruction") or "").strip() and gid in by_group:
            q["instruction"] = by_group[gid]
    return questions


def _group_consecutive_gap_fill(questions: list[dict]) -> list[dict]:
    """Kitobdagi 'Complete the sentences' kabi mashqlarni yaxlit qiladi.

    - AI `group` id bergan gap_fill savollar bitta passage_cloze ga birlashadi.
    - Aks holda ketma-ket 3+ ta gap_fill ham bitta passage_cloze ga birlashadi.
    Boshqa turlar (matching, listening, word_practice...) allaqachon yaxlit — tegilmaydi.
    """
    if not questions:
        return questions
    grouped_ids: dict[str, list[dict]] = {}
    for q in questions:
        gid = q.get("group") if q.get("kind") == "gap_fill" else None
        if gid:
            grouped_ids.setdefault(str(gid), []).append(q)
    merge_ids = {gid for gid, r in grouped_ids.items() if len(r) >= 2}

    out: list[dict] = []
    consumed_ids: set[str] = set()
    run: list[dict] = []

    def _flush():
        nonlocal run
        if len(run) >= 3:
            out.append(_cloze_from_gap_fill_run(run))
        else:
            out.extend(run)
        run = []

    for q in questions:
        gid = q.get("group") if q.get("kind") == "gap_fill" else None
        if gid and str(gid) in merge_ids:
            _flush()
            if str(gid) not in consumed_ids:
                consumed_ids.add(str(gid))
                out.append(_cloze_from_gap_fill_run(grouped_ids[str(gid)]))
            continue
        if q.get("kind") == "gap_fill" and str(q.get("answer") or "").strip():
            run.append(q)
        else:
            _flush()
            out.append(q)
    _flush()
    return out


def _materialize_word_practice(q: dict, lang: str = "Uzbek", study_lang: str = "English") -> dict:
    """word_practice ni random konkret test turiga aylantiradi (har attemptda boshqacha).

    Ko'rsatma (prompt) o'rganilayotgan til (study_lang: English/Russian) da yoziladi.
    Tarjima esa student guruh tiliga (lang: Uzbek/Russian) qarab so'raladi."""
    import random

    ru = lang == "Russian"
    study_ru = study_lang == "Russian"
    word = str(q.get("word") or q.get("prompt") or "").strip()
    # Guruh tiliga mos tarjima: ruscha guruh -> ruscha, aks holda o'zbekcha.
    tr_uz = str(q.get("translation_uz") or "").strip()
    tr_ru = str(q.get("translation_ru") or "").strip()
    legacy = str(q.get("translation") or "").strip()
    translation = (tr_ru if ru else tr_uz) or legacy or tr_uz or tr_ru
    # Ikkala tarjima ham qabul qilinadi (student boshqa tilda yozsa ham).
    accepted = [t for t in {translation, tr_uz, tr_ru, legacy} if t]
    instruction = q.get("instruction")
    level = q.get("level")
    variants = list(WORD_PRACTICE_VARIANTS)
    if not translation and "translation" in variants:
        variants.remove("translation")
    kind = random.choice(variants) if variants else "write_sentence"
    base: dict[str, Any] = {"kind": kind, "word": word, "level": level, "instruction": instruction}
    if kind == "speak_sentence":
        base["prompt"] = q.get("prompt") or (
            f"Составьте и произнесите предложение со словом «{word}»" if study_ru
            else f"Make and say a sentence using '{word}'"
        )
    elif kind == "write_sentence":
        base["prompt"] = q.get("prompt") or (
            f"Напишите предложение со словом «{word}»" if study_ru
            else f"Write a sentence using '{word}'"
        )
        base["reference_answer"] = None
    elif kind == "spelling":
        base["prompt"] = q.get("prompt") or (
            "Правильно напишите слово" if study_ru
            else "Spell the word correctly"
        )
        base["answer"] = word
        base["accepted_answers"] = []
    else:  # translation — tarjima student tiliga (uz/ru)
        base["prompt"] = (f"Переведите: {word}" if ru else f"Tarjima qiling: {word}")
        base["answer"] = translation
        base["accepted_answers"] = accepted
    return base


def _expand_polymorphic_questions(questions: list[dict], lang: str = "Uzbek", study_lang: str = "English") -> list[dict]:
    """Attempt boshlanishidan oldin polimorf savollarni konkret turga ochadi."""
    out: list[dict] = []
    for q in questions or []:
        if str(q.get("kind") or "") in POLYMORPHIC_KINDS:
            out.append(_materialize_word_practice(q, lang, study_lang))
        else:
            out.append(q)
    return out


def _question_for_student(question: dict) -> dict:
    """Javobni yashirib, studentga ko'rsatiladigan shaklga keltiradi."""
    kind = str(question.get("kind") or "")
    meta = AI_TEST_TYPES.get(kind, {})
    out = {
        "kind": kind,
        "check": meta.get("check", "ai"),
        "input": meta.get("input", "text"),
        "retry_until_correct": bool(meta.get("retry_until_correct", True)),
        "prompt": question.get("prompt"),
        "instruction": question.get("instruction"),
        "word": question.get("word"),
        "passage": question.get("passage"),
        "image_url": question.get("image_url"),
        "audio_url": question.get("audio_url"),
        "level": question.get("level"),
    }
    if kind == "listening":
        out["options"] = question.get("options") or []
    elif kind == "matching":
        pairs = question.get("pairs") or []
        out["left_items"] = [p.get("left") for p in pairs]
        out["right_items"] = sorted([p.get("right") for p in pairs])
    elif kind == "scrambled_sentence":
        # To'g'ri so'zlar + AI chalg'ituvchi so'zlari aralashtiriladi.
        tokens = list(question.get("tokens") or []) + list(question.get("distractors") or [])
        random.shuffle(tokens)
        out["tokens"] = tokens
    elif kind == "passage_cloze":
        # Yaxlit: matn (___ bilan), bo'sh joylar soni va so'zlar banki.
        out["passage_template"] = question.get("passage_template") or question.get("passage") or ""
        out["blank_count"] = len(question.get("blanks") or [])
        out["word_bank"] = list(question.get("word_bank") or [])
    elif kind == "reading_set":
        # Matn + savollar (javoblarsiz) bitta kartada.
        out["passage"] = question.get("passage") or ""
        out["sub_questions"] = [
            {
                "type": s.get("type"),
                "prompt": s.get("prompt"),
                "options": list(s.get("options") or []),
            }
            for s in (question.get("sub_questions") or [])
        ]
    return out


# ═══════════════════════════════════════════════════════════════════════════
# 2. AI SKRINSHOT IMPORT — kitob sahifasidan matn + mashqlar
# ═══════════════════════════════════════════════════════════════════════════

class ScreenshotImportRequest(BaseModel):
    # Har qanday fayl turi: rasm (png/jpg/webp…), PDF, DOC/DOCX, TXT/RTF.
    # `image_urls` — eski mijozlar bilan moslik uchun; `file_urls` — yangi nom.
    image_urls: list[str] = Field(default_factory=list, max_length=10)
    file_urls: list[str] = Field(default_factory=list, max_length=10)
    subject: str = "English"
    level: str | None = None
    node_id: int | None = None
    instruction: str | None = None
    #: 'uz' yoki 'ru' — yevro/rus guruh uchun 'ru' (tarjima+ko'rsatma ruscha)
    lang: str | None = None
    #: bo'sh bo'lsa AI rasmda ko'rgan mashq turlarining o'zini tanlaydi
    wanted_kinds: list[str] | None = None

    @property
    def all_urls(self) -> list[str]:
        seen: list[str] = []
        for url in [*(self.file_urls or []), *(self.image_urls or [])]:
            clean = str(url or "").strip()
            if clean and clean not in seen:
                seen.append(clean)
        return seen


def _import_system_prompt(subject: str, level: str | None, wanted: list[str], instruction_language: str) -> str:
    kinds_help = "\n".join(
        f'- "{kind}": {meta["label_en"]}'
        + (" (teacher must upload the audio file afterwards)" if meta["needs_audio_asset"] else "")
        + (" — polymorphic: for a single vocabulary word; the app randomly turns it into a "
           "speaking/writing/spelling/translation task when the student reaches it" if meta.get("polymorphic") else "")
        for kind, meta in AI_TEST_TYPES.items()
        if kind in wanted
    )
    return (
        "You are an expert curriculum digitizer for a language school.\n"
        "You receive coursebook material — as page images and/or extracted text from PDF/DOC/DOCX.\n\n"
        "=== VOCABULARY LIST RULES (highest priority) ===\n"
        "• A vocabulary list is lines of the form: word — translation or word (pos) — translation.\n"
        "• Extract EVERY word that is LITERALLY listed. A 50-word list → exactly 50 word_practice items.\n"
        "• Do NOT add, invent, infer, or generate ANY word that is not physically written in the list.\n"
        "• Do NOT add synonyms, related words, or words from example sentences.\n"
        "• Each word_practice item must include translation_uz (Uzbek) and translation_ru (Russian) from the list.\n\n"
        "=== READING TEXT RULES ===\n"
        "• Whenever the material contains a reading passage, story, dialogue or article:\n"
        "  1. Put the COMPLETE verbatim text in reading_text field.\n"
        "  2. Output ONE reading_set question: put full text in 'passage', generate EXACTLY 5 comprehension\n"
        "     questions in 'questions'. Use a MIX of these types:\n"
        "     - true_false_ng (True/False/Not Given — 3 options)\n"
        "     - choice (multiple choice with 4 options, answer verifiable from text)\n"
        "     - synonym (find a word in the text meaning X)\n"
        "     - gap (a sentence from the text with one blank ___)\n"
        "     - who_said (who said/did something in the text)\n"
        "  3. ALL answers MUST be verifiable directly from the text — no inference.\n"
        "  4. Do this even if the original material has NO printed questions for the text.\n\n"
        "=== PRINTED EXERCISE RULES ===\n"
        "0. Analyze page structure first. A single printed exercise (one instruction + its items,\n"
        "   e.g. 'Complete the sentences. Use a verb from the box' + 8 sentences, OR a reading\n"
        "   text with numbered gaps) = ONE cohesive task → output as ONE 'passage_cloze'.\n"
        "   Give every question a 'group' string = the exercise number (e.g. '11.1').\n"
        "   Only use standalone 'gap_fill' for a truly independent single gap.\n"
        "1. Convert ONLY what is LITERALLY printed. Do NOT invent exercises not on the page.\n"
        "1b. NUMBERED exercises: keep order and numbering. For a numbered gap list sharing one\n"
        "   instruction/word box → ONE 'passage_cloze' with each item on its OWN line prefixed\n"
        "   by its number ('1. get — ___\\n2. see — ___'). Continuous prose stays as flowing text.\n"
        "2. Use the question kind that matches each printed exercise exactly.\n"
        "2b. Every question MUST have its own 'kind' field (exact snake_case from allowed list).\n"
        "3. Keep target-language content exactly as written. Don't translate unless the exercise is a translation task.\n"
        "4. For every auto-checked type provide the exact expected answer.\n"
        "4b. For 'scrambled_sentence' add 'distractors': 2–4 extra plausible words not in the correct sentence.\n"
        "5. Listening exercises: produce the question, leave audio_url empty (teacher uploads audio later).\n"
        "6. LANGUAGE OF INSTRUCTIONS: copy the book's own wording verbatim. English coursebook → English\n"
        "   instructions. Only translated fields are word_practice's translation_uz/translation_ru.\n"
        "6b. reading_text MUST be COMPLETE, VERBATIM text of ALL reading passages/dialogues.\n"
        "6c. WORKBOOK EXERCISE FORMATS:\n"
        "   • 'Complete the sentences' with WORD BANK → ONE 'passage_cloze', each sentence on its OWN line.\n"
        "   • Continuous passage with numbered gaps → ONE 'passage_cloze', flowing prose.\n"
        "   • Single-sentence gap with no shared text → 'gap_fill'.\n"
        "   • Verb/word transformation lists → ONE 'passage_cloze', each item on its own line.\n"
        "   • 'Write sentences about...' → ONE 'write_sentence' PER numbered item.\n"
        "7. If the material is not educational, return {\"error\":\"not_educational\"}.\n\n"
        f"Subject: {subject}. Level: {level or 'infer from the material'}. "
        f"Translations language for word_practice: Uzbek + Russian (provide both).\n\n"
        f"Allowed question kinds:\n{kinds_help}\n\n"
        "Return ONLY valid JSON, no markdown fences, with this shape:\n"
        "{\n"
        '  "title": "short title of the page/unit",\n'
        '  "level": "A1|A2|B1|B2|C1|C2 or empty",\n'
        '  "reading_text": "the main text of the page (full verbatim), empty string if none",\n'
        '  "notes": "grammar rule/explanation found on the page, empty if none",\n'
        '  "questions": [\n'
        '    {"kind":"word_practice","word":"decide","translation_uz":"qaror qilmoq","translation_ru":"решать"},\n'
        '    {"kind":"gap_fill","prompt":"She ___ to school every day.","answer":"goes",'
        '"accepted_answers":["goes"],"instruction":"Complete the sentence."},\n'
        '    {"kind":"matching","prompt":"Match the words to the definitions.",'
        '"pairs":[{"left":"brave","right":"not afraid"}]},\n'
        '    {"kind":"scrambled_sentence","prompt":"Put the words in order.",'
        '"answer":"I have never been to Paris.","tokens":["I","have","never","been","to","Paris."],'
        '"distractors":["was","going","the"]},\n'
        '    {"kind":"reading_set","passage":"Alex was crazy about sport and music...",'
        '"questions":[{"type":"true_false_ng","prompt":"Alex was a member of a local club.","answer":"True",'
        '"options":["True","False","Not given"]},'
        '{"type":"choice","prompt":"What was Alex interested in?","options":["Sport and music","Only chess","Art","Dancing"],"answer":"Sport and music"},'
        '{"type":"synonym","prompt":"Find a word meaning \'very enthusiastic about\'.","answer":"crazy about"},'
        '{"type":"gap","prompt":"He was a member of a local ___.","answer":"club"},'
        '{"type":"who_said","prompt":"Who trained in the gym with his team?","answer":"Alex"}]}\n'
        '    {"kind":"passage_cloze","instruction":"Put the verbs in the correct form.",'
        '"passage":"Last Tuesday Lisa ___ from London to Madrid.",'
        '"answers":[{"answer":"flew"}],"word_bank":["fly","get","have"]}\n'
        "  ]\n"
        "}\n"
        "Produce as many questions as the material contains. "
        "For vocabulary lists: one word_practice per listed word, no extras. "
        "For reading texts: always include a reading_set with 5 comprehension questions."
    )


def _instruction_language_name(code: str | None, subject: str) -> str:
    c = str(code or "").strip().lower()[:2]
    if c == "ru":
        return "Russian"
    if c == "uz":
        return "Uzbek"
    # Fan rus tili bo'lsa — rus, aks holda o'zbek (default).
    subj = str(subject or "").lower()
    if any(t in subj for t in ("rus", "рус", "russian")):
        return "Russian"
    return "Uzbek"


def _study_language_name(subject: str) -> str:
    """O'rganilayotgan til — mashq ko'rsatmalari shu tilda bo'ladi.
    English kursi -> English, Russian kursi -> Russian, aks holda English."""
    subj = str(subject or "").lower()
    if any(t in subj for t in ("rus", "рус", "russian")):
        return "Russian"
    return "English"


@router.post("/teacher/library/ai/import-screenshot")
async def library_ai_import_screenshot(
    payload: ScreenshotImportRequest,
    authorization: str | None = Header(default=None),
    x_language: str | None = Header(default=None, alias="X-Language"),
):
    """Har qanday fayldan (rasm/PDF/DOC/DOCX/TXT) fayldagi BARCHA o'quv materialini
    har xil, random test turlariga aylantiradi. Ko'rsatmalar o'quv tilida (rus/o'zbek)
    avtomatik yoziladi."""
    user = _auth(authorization, TEACHER_ROLES)
    if payload.node_id:
        node = _node_or_404(int(payload.node_id))
        _require_node_permission(user, node, "edit")

    urls = payload.all_urls
    if not urls:
        raise HTTPException(status_code=400, detail="Fayl yuborilmadi")

    wanted = [k for k in (payload.wanted_kinds or []) if k in AI_TEST_TYPES] or list(AI_TEST_TYPES.keys())
    # Til: aniq tanlangan (yevro guruh -> 'ru') > X-Language > fandan avtomatik.
    instruction_language = _instruction_language_name(payload.lang or x_language, payload.subject)
    vocab_lang = "ru" if instruction_language == "Russian" else "uz"
    vision_urls, text_blocks, unsupported = _prepare_import_sources(urls, owner_id=int(user.get("id") or 0))
    if not vision_urls and not text_blocks:
        detail = "Fayldan matn yoki rasm o'qib bo'lmadi"
        if unsupported:
            detail += f" ({', '.join(unsupported)} qo'llab-quvvatlanmaydi)"
        raise HTTPException(status_code=422, detail=detail)

    from ai_generator import _xai_generate_text, _xai_generate_text_stream_with_images
    import aiohttp

    user_prompt = str(payload.instruction or "").strip() or (
        "Digitize this coursebook material: extract its text and convert every exercise into structured questions."
    )
    system_prompt = _import_system_prompt(payload.subject, payload.level, wanted, instruction_language)
    extracted = "\n\n".join(text_blocks).strip()
    if extracted:
        # 60k belgidan oshsa kesamiz (juda katta hujjatlarni AI baribir hazm qilmaydi).
        extracted = extracted[:60000]

    raw = ""

    async def _run_ai(session) -> str:
        if vision_urls:
            task = system_prompt + "\n\nTASK: " + user_prompt
            if extracted:
                task += "\n\nAdditional extracted text from the same material:\n" + extracted
            chunks: list[str] = []
            async for chunk in _xai_generate_text_stream_with_images(task, vision_urls, session=session):
                chunks.append(str(chunk or ""))
            return "".join(chunks).strip()
        task = (
            system_prompt
            + "\n\nTASK: " + user_prompt
            + "\n\nMATERIAL TEXT:\n" + extracted
        )
        return await _xai_generate_text(
            task,
            session=session,
            system_content="You are an educational content digitizer. Output only valid JSON.",
        )

    # Birinchi urinishda AI ba'zan bo'sh/qisqa javob qaytaradi — 2 marta urinamiz.
    data: dict = {}
    questions: list[dict] = []
    last_error: Exception | None = None
    try:
        async with aiohttp.ClientSession() as session:
            for attempt_no in range(2):
                try:
                    raw = await _run_ai(session)
                except Exception as exc:
                    last_error = exc
                    logger.warning("library ai import attempt %s failed: %s", attempt_no + 1, exc)
                    continue
                logger.debug(
                    "library ai import raw attempt=%s chars=%d preview=%s",
                    attempt_no + 1, len(raw or ""), (raw or "")[:300].replace("\n", " ")
                )
                data = _extract_json_object(str(raw or "").strip())
                logger.debug(
                    "library ai import parsed attempt=%s data_keys=%s questions_raw_len=%s",
                    attempt_no + 1, list(data.keys()), len(data.get("questions") or [])
                )
                if str(data.get("error") or "") == "not_educational":
                    raise HTTPException(status_code=422, detail="Fayl o'quv materiali emas")
                questions = _normalize_questions(data.get("questions"))
                logger.info(
                    "library ai import normalized attempt=%s questions=%d",
                    attempt_no + 1, len(questions)
                )
                if questions:
                    break  # muvaffaqiyat
    except HTTPException:
        raise
    except Exception as exc:
        logger.exception("library ai import failed teacher=%s: %s", user.get("id"), exc)
        last_error = exc

    # Lug'at (vocabulary) ro'yxatini deterministik ravishda ajratamiz —
    # FAQAT AI nol savol qaytargan bo'lsa (timeout/xato/bo'sh javob) fallback sifatida.
    ai_word_count = sum(1 for q in questions if q.get("kind") == "word_practice")
    if text_blocks and ai_word_count == 0:
        # AI so'zlarni ajrata olmadi — deterministik parser bilan to'ldiramiz.
        existing_words = {
            str(q.get("word") or "").strip().lower()
            for q in questions if q.get("kind") == "word_practice"
        }
        for item in _extract_vocab_items("\n".join(text_blocks), vocab_lang):
            key = item["word"].strip().lower()
            if key and key not in existing_words:
                questions.append(item)
                existing_words.add(key)
        logger.info(
            "library ai import vocab fallback used: added=%d total_questions=%d",
            sum(1 for q in questions if q.get("kind") == "word_practice"),
            len(questions)
        )

    if not questions and last_error is not None:
        raise HTTPException(
            status_code=503,
            detail=f"AI hozir javob bermadi ({str(last_error)[:120]}). Qayta urinib ko'ring.",
        )

    # Ketma-ket mayda gap_fill savollarni yaxlit passage_cloze mashqiga birlashtiramiz
    # (AI alohida chiqarsa ham — kitobdagidek bitta yaxlit karta bo'ladi).
    questions = _propagate_group_instructions(questions)
    questions = _group_consecutive_gap_fill(questions)
    if not questions:
        raise HTTPException(status_code=422, detail="Fayldan mashq topilmadi. Aniqroq material yuboring.")
    needs_audio = [
        {"index": i, "kind": q["kind"], "prompt": q.get("prompt")}
        for i, q in enumerate(questions)
        if q.get("needs_audio_upload")
    ]
    return {
        "message": f"{len(questions)} ta mashq tayyorlandi",
        "title": str(data.get("title") or "").strip() or "Yangi mavzu",
        "level": str(data.get("level") or payload.level or "").strip() or None,
        "reading_text": str(data.get("reading_text") or "").strip(),
        "notes": str(data.get("notes") or "").strip(),
        "questions": questions,
        "needs_audio_upload": needs_audio,
        "kinds_used": sorted({q["kind"] for q in questions}),
        "unsupported": unsupported,
    }


_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp", ".avif", ".ico", ".tiff", ".tif"}
_PDF_EXTS = {".pdf"}
_DOC_EXTS = {".docx", ".doc", ".odt", ".ott"}
_SHEET_EXTS = {".xlsx", ".xls", ".ods", ".csv"}
_SLIDE_EXTS = {".pptx", ".ppt", ".odp", ".otp"}
_TEXT_EXTS = {".txt", ".rtf", ".md", ".tex", ".log", ".xml", ".html", ".htm"}
_ALL_KNOWN_EXTS = _IMAGE_EXTS | _PDF_EXTS | _DOC_EXTS | _SHEET_EXTS | _SLIDE_EXTS | _TEXT_EXTS
_MAX_VISION_IMAGES = 3
_MAX_PDF_PAGES = 3

# Lug'at qatori: "headword (pos) — tarjima — перевод" ko'rinishidagi qatorlarni
# ushlaydigan ajratgichlar (em/en dash yoki oddiy tire).
_VOCAB_SEPARATORS = ["—", "–", " - ", " – ", " — "]


def _extract_vocab_items(text: str, lang: str = "uz") -> list[dict]:
    """Matndan lug'at ro'yxatini to'liq ajratib, word_practice mashqlariga aylantiradi.

    Har bir qator: `beat (v) — mag'lub etmoq — побеждать` ko'rinishida bo'ladi.
    Har bir so'z uchun HAM o'zbekcha, HAM ruscha tarjima ajratiladi (bittada) —
    guruh tiliga qarab alohida test qilinmaydi. Nasr qatorlari (ajratgichsiz) tushadi.
    """
    def _is_cyrillic(s: str) -> bool:
        return bool(re.search(r"[А-Яа-яЁё]", s))

    items: list[dict] = []
    seen: set[str] = set()
    for raw_line in str(text or "").splitlines():
        line = raw_line.strip()
        if not line or len(line) > 200:
            continue
        sep = next((s for s in _VOCAB_SEPARATORS if s in line), None)
        if not sep:
            continue
        parts = [p.strip() for p in line.split(sep) if p.strip()]
        if len(parts) < 2:
            continue
        head = parts[0]
        # Sarlavha (masalan "Elementary: Unit 3") lug'at so'zi emas — o'tkazamiz.
        if ":" in head:
            continue
        # Sarlavha/nasr emas: chap tomon qisqa (so'z yoki ibora) bo'lishi kerak.
        head_wordcount = len(re.sub(r"\([^)]*\)", "", head).split())
        if head_wordcount == 0 or head_wordcount > 5:
            continue
        if not re.search(r"[A-Za-zА-Яа-яЁё]", head):
            continue
        # So'z (pos belgisisiz).
        word = re.sub(r"\s*\([^)]*\)\s*", " ", head).strip(" .:;")
        if not word or word.lower() in seen:
            continue
        glosses = [g.strip(" .:;") for g in parts[1:] if g.strip()]
        if not glosses:
            continue
        # Ruscha (kirill) va o'zbekcha (lotin) glossalarni ajratamiz.
        ru_gloss = next((g for g in glosses if _is_cyrillic(g)), None)
        uz_gloss = next((g for g in glosses if not _is_cyrillic(g)), None)
        # Ikkalasi ham juda uzun bo'lsa — bu nasr, o'tkazamiz.
        candidate = uz_gloss or ru_gloss or ""
        if len(candidate.split()) > 10:
            continue
        seen.add(word.lower())
        items.append({
            "kind": "word_practice",
            "word": word,
            "translation_uz": uz_gloss or None,
            "translation_ru": ru_gloss or None,
            # Backward-compat: eski maydonlar ham to'ldiriladi.
            "translation": uz_gloss or ru_gloss or None,
            "meaning": ru_gloss if uz_gloss else None,
        })
    return items


def _resolve_local_path(url: str):
    """Yuklangan fayl serverdagi haqiqiy yo'lini topadi (bir nechta upload papka)."""
    from backend.main import HOMEWORK_UPLOAD_DIR

    name = Path(urlparse(str(url or "")).path).name
    if not name:
        return None
    candidates = [
        HOMEWORK_UPLOAD_DIR / name,
        HOMEWORK_UPLOAD_DIR.parent / "homework_uploads" / name,
    ]
    for cand in candidates:
        try:
            if cand.is_file():
                return cand
        except Exception:
            continue
    return None


def _abs_url(url: str) -> str:
    base = str(os.getenv("BACKEND_BASE_URL") or "https://diamond-education.uz/api").rstrip("/")
    clean = str(url or "").strip()
    return base + clean if clean.startswith("/") else clean


def _prepare_import_sources(urls: list[str], owner_id: int) -> tuple[list[str], list[str], list[str]]:
    """Har xil fayllarni AI uchun tayyorlaydi (universal — har qanday fayl qabul qilinadi).

    Qaytaradi: (vision uchun rasm URL lari, ajratilgan matn bloklari, qo'llab-
    quvvatlanmaydigan kengaytmalar ro'yxati). PDF sahifalari PNG ga render
    qilinib vision ro'yxatiga qo'shiladi; hujjatlar matn sifatida o'qiladi;
    noma'lum kengaytmalar uchun universal fallback ishlatiladi.
    """
    vision_urls: list[str] = []
    text_blocks: list[str] = []
    unsupported: list[str] = []

    for url in urls:
        ext = Path(urlparse(str(url or "")).path).suffix.lower()
        logger.debug("import_sources: url=%s ext=%s", url, ext)

        # 1. Rasm — to'g'ridan vision API ga
        if ext in _IMAGE_EXTS:
            if len(vision_urls) < _MAX_VISION_IMAGES:
                vision_urls.append(_abs_url(url))
            continue

        # Barcha boshqa turlar uchun lokal faylni topamiz
        local = _resolve_local_path(url)
        if local is None:
            logger.warning("import_sources: local file not found for url=%s", url)
            unsupported.append(ext or "unknown")
            continue

        # 2. PDF — avval render, bo'lmasa matn
        if ext in _PDF_EXTS:
            rendered = _render_pdf_to_images(local, owner_id)
            if rendered:
                for served in rendered:
                    if len(vision_urls) < _MAX_VISION_IMAGES:
                        vision_urls.append(_abs_url(served))
            else:
                text = _extract_pdf_text(local)
                if text.strip():
                    text_blocks.append(text)
                    logger.info("import_sources: pdf->text chars=%d url=%s", len(text), url)
                else:
                    unsupported.append(ext)
            continue

        # 3. Word/DOC/DOCX
        if ext in _DOC_EXTS:
            text = _extract_docx_text(local)
            if text.strip():
                text_blocks.append(text)
                logger.info("import_sources: docx->text chars=%d url=%s", len(text), url)
            else:
                unsupported.append(ext)
            continue

        # 4. Excel/Spreadsheet
        if ext in _SHEET_EXTS:
            text = _extract_xlsx_text(local)
            if text.strip():
                text_blocks.append(text)
                logger.info("import_sources: sheet->text chars=%d url=%s", len(text), url)
            else:
                unsupported.append(ext)
            continue

        # 5. PowerPoint/Presentation
        if ext in _SLIDE_EXTS:
            text = _extract_pptx_text(local)
            if text.strip():
                text_blocks.append(text)
                logger.info("import_sources: pptx->text chars=%d url=%s", len(text), url)
            else:
                unsupported.append(ext)
            continue

        # 6. Matn/HTML/XML formatlar
        if ext in _TEXT_EXTS:
            if ext in (".html", ".htm", ".xml"):
                text = _extract_html_text(local)
            else:
                text = _read_text_file(local)
            if text.strip():
                text_blocks.append(text)
                logger.info("import_sources: text->text chars=%d url=%s", len(text), url)
            else:
                unsupported.append(ext)
            continue

        # 7. Noma'lum kengaytma — universal fallback
        logger.info("import_sources: unknown ext=%s, trying universal extract url=%s", ext, url)
        text = _extract_any_text(local)
        if text.strip():
            text_blocks.append(text)
            logger.info("import_sources: universal->text chars=%d url=%s", len(text), url)
        else:
            # Oxirgi urinish: vision API ga rasm sifatida yuboramiz (agar kichik bo'lsa)
            fsize = local.stat().st_size if local.exists() else 0
            if fsize < 10 * 1024 * 1024 and len(vision_urls) < _MAX_VISION_IMAGES:
                vision_urls.append(_abs_url(url))
                logger.info("import_sources: unknown->vision url=%s", url)
            else:
                unsupported.append(ext or "unknown")

    return vision_urls, text_blocks, unsupported


def _render_pdf_to_images(path, owner_id: int) -> list[str]:
    """PDF sahifalarini PNG ga aylantiradi va servisdan yuklanadigan URL qaytaradi."""
    try:
        import fitz  # PyMuPDF
    except Exception:
        logger.info("PyMuPDF yo'q — PDF matn sifatida o'qiladi")
        return []
    from backend.main import HOMEWORK_UPLOAD_DIR

    served: list[str] = []
    try:
        doc = fitz.open(str(path))
        try:
            page_count = min(_MAX_PDF_PAGES, doc.page_count)
            for i in range(page_count):
                page = doc.load_page(i)
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # ~144 DPI
                filename = f"library_pdf_{owner_id}_{int(__import__('time').time()*1000)}_{i}.png"
                out_path = HOMEWORK_UPLOAD_DIR / filename
                pix.save(str(out_path))
                served.append(f"/homework/files/{filename}")
        finally:
            doc.close()
    except Exception:
        logger.exception("PDF render failed path=%s", path)
        return []
    return served


def _extract_pdf_text(path) -> str:
    try:
        import fitz

        doc = fitz.open(str(path))
        try:
            parts = [doc.load_page(i).get_text() for i in range(min(20, doc.page_count))]
            return "\n".join(p for p in parts if p)
        finally:
            doc.close()
    except Exception:
        return ""


def _extract_docx_text(path) -> str:
    try:
        import docx  # python-docx

        document = docx.Document(str(path))
        lines = [p.text for p in document.paragraphs if p.text and p.text.strip()]
        for table in document.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text and c.text.strip()]
                if cells:
                    lines.append(" | ".join(cells))
        return "\n".join(lines)
    except Exception:
        logger.exception("docx extract failed path=%s", path)
        return ""


def _extract_xlsx_text(path) -> str:
    """Excel (.xlsx/.xls/.ods) fayldan barcha sahifa matnini ajratadi."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
        lines: list[str] = []
        for ws in wb.worksheets:
            ws_lines: list[str] = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
                if cells:
                    ws_lines.append(" | ".join(cells))
            if ws_lines:
                lines.append(f"=== {ws.title} ===\n" + "\n".join(ws_lines))
        wb.close()
        return "\n\n".join(lines)
    except Exception:
        pass
    # CSV fallback
    try:
        import csv as _csv
        raw = Path(path).read_bytes()
        for enc in ("utf-8", "cp1251", "latin-1"):
            try:
                text = raw.decode(enc)
                rows = list(_csv.reader(text.splitlines()))
                return "\n".join(" | ".join(r) for r in rows if any(c.strip() for c in r))
            except Exception:
                continue
    except Exception:
        pass
    logger.warning("xlsx/xls extract failed path=%s", path)
    return ""


def _extract_pptx_text(path) -> str:
    """PowerPoint (.pptx/.odp) fayldan barcha slayd matnini ajratadi."""
    # python-pptx urinishi
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        lines: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            slide_lines: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_lines.append(shape.text.strip())
            if slide_lines:
                lines.append(f"--- Slide {i} ---\n" + "\n".join(slide_lines))
        return "\n\n".join(lines)
    except Exception:
        pass
    # zipfile + XML fallback (python-pptx yo'q bo'lsa)
    try:
        import zipfile, re as _re
        lines: list[str] = []
        with zipfile.ZipFile(str(path)) as z:
            slide_names = sorted(n for n in z.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml"))
            for sname in slide_names:
                xml = z.read(sname).decode("utf-8", errors="ignore")
                texts = _re.findall(r"<a:t>([^<]+)</a:t>", xml)
                if texts:
                    lines.append(" ".join(t.strip() for t in texts if t.strip()))
        return "\n".join(lines)
    except Exception:
        logger.warning("pptx extract failed path=%s", path)
        return ""


def _extract_odt_text(path) -> str:
    """ODF (ODT/ODS/ODP) fayldan matn ajratadi (zipfile + XML)."""
    try:
        import zipfile, re as _re
        with zipfile.ZipFile(str(path)) as z:
            if "content.xml" not in z.namelist():
                return ""
            xml = z.read("content.xml").decode("utf-8", errors="ignore")
        # <text:p> va <text:h> teglaridan matnni ajratamiz
        texts = _re.findall(r"<text:[ph][^>]*>([^<]*(?:<[^/][^>]*>[^<]*</[^>]*>)*[^<]*)</text:[ph]>", xml)
        if not texts:
            texts = _re.findall(r">([^<]{2,})<", xml)
        lines = [t.strip() for t in texts if t.strip() and len(t.strip()) > 1]
        return "\n".join(lines)
    except Exception:
        logger.warning("odt extract failed path=%s", path)
        return ""


def _extract_html_text(path) -> str:
    """HTML/XML fayldan oddiy matnni ajratadi."""
    try:
        import re as _re
        raw = Path(path).read_bytes()
        for enc in ("utf-8", "cp1251", "latin-1"):
            try:
                text = raw.decode(enc)
                break
            except Exception:
                continue
        else:
            text = raw.decode("utf-8", errors="ignore")
        # Script va style teglarini olib tashlaymiz
        text = _re.sub(r"<(script|style)[^>]*>.*?</(script|style)>", " ", text, flags=_re.DOTALL | _re.IGNORECASE)
        # Barcha teglarni olib tashlaymiz
        text = _re.sub(r"<[^>]+>", " ", text)
        # Encoding entities
        text = text.replace("&amp;", "&").replace("&lt;", "<").replace("&gt;", ">").replace("&nbsp;", " ")
        # Bo'sh qatorlarni tozalaymiz
        lines = [l.strip() for l in text.splitlines() if l.strip()]
        return "\n".join(lines)
    except Exception:
        logger.warning("html extract failed path=%s", path)
        return ""


def _extract_any_text(path) -> str:
    """Har qanday fayldan matn ajratishga urinadi (universal fallback)."""
    if path is None:
        return ""
    p = Path(str(path))
    ext = p.suffix.lower()

    # Tartiblangan urinishlar
    if ext in _DOC_EXTS:
        return _extract_docx_text(p)
    if ext in _SHEET_EXTS:
        return _extract_xlsx_text(p)
    if ext in _SLIDE_EXTS:
        return _extract_pptx_text(p)
    if ext in {".odt", ".ott", ".ods", ".odp"}:
        return _extract_odt_text(p)
    if ext in _TEXT_EXTS:
        return _read_text_file(p)

    # Noma'lum kengaytma — bosqichma-bosqich sinab ko'ramiz
    # 1. ZIP asosli (docx/xlsx/pptx/odt hammasi ZIP)
    try:
        import zipfile
        if zipfile.is_zipfile(str(p)):
            # DOCX urinishi
            t = _extract_docx_text(p)
            if t.strip():
                return t
            # XLSX urinishi
            t = _extract_xlsx_text(p)
            if t.strip():
                return t
            # PPTX urinishi
            t = _extract_pptx_text(p)
            if t.strip():
                return t
            # ODT urinishi
            t = _extract_odt_text(p)
            if t.strip():
                return t
    except Exception:
        pass

    # 2. PDF urinishi
    try:
        t = _extract_pdf_text(p)
        if t.strip():
            return t
    except Exception:
        pass

    # 3. Oddiy matn sifatida o'qish (UTF-8 / CP1251 / latin)
    try:
        t = _read_text_file(p)
        if t.strip():
            return t
    except Exception:
        pass

    logger.warning("universal extract failed path=%s ext=%s", path, ext)
    return ""

def _read_text_file(path) -> str:
    try:
        raw = Path(path).read_bytes()
        for enc in ("utf-8", "cp1251", "latin-1"):
            try:
                return raw.decode(enc)
            except Exception:
                continue
        return raw.decode("utf-8", errors="ignore")
    except Exception:
        return ""


def _extract_json_object(raw: str) -> dict:
    text = str(raw or "").strip()
    if text.startswith("```"):
        text = re.sub(r"^```[a-zA-Z]*\s*", "", text)
        text = re.sub(r"```\s*$", "", text).strip()
    try:
        parsed = json.loads(text)
        if isinstance(parsed, dict):
            return parsed
    except Exception:
        pass
    match = re.search(r"\{.*\}", text, re.DOTALL)
    if match:
        try:
            parsed = json.loads(match.group(0))
            return parsed if isinstance(parsed, dict) else {}
        except Exception:
            return {}
    return {}


# ═══════════════════════════════════════════════════════════════════════════
# 3. STUDENT AI TEST RUNTIME
#    - taymer yo'q
#    - har bir javob darhol tekshiriladi
#    - xato bo'lsa shu savolda qolinadi (retry_until_correct)
#    - testdan chiqib ketilsa attempt bekor bo'ladi
# ═══════════════════════════════════════════════════════════════════════════

class AiTestStartRequest(BaseModel):
    source_type: str = "library_test"
    source_id: int | None = None
    homework_id: int | None = None


class AiTestAnswerRequest(BaseModel):
    question_index: int = Field(ge=0)
    answer_text: str | None = None
    audio_url: str | None = None
    #: matching uchun {"left": "right"}, scrambled uchun so'zlar ketma-ketligi
    pairs: dict[str, str] | None = None
    order: list[str] | None = None
    choice_index: int | None = None
    #: passage_cloze uchun bo'sh joylar javoblari (tartib bo'yicha)
    blanks: list[str] | None = None


def _attempt_state(attempt: dict) -> dict:
    """Attemptni studentga ko'rsatiladigan holatga aylantiradi."""
    questions = attempt.get("questions") or []
    answers = attempt.get("answers") or []
    solved: set[int] = set()
    tries: dict[int, int] = {}
    feedback: dict[int, dict] = {}
    for row in answers:
        idx = int(row.get("question_index") or 0)
        tries[idx] = max(tries.get(idx, 0), int(row.get("try_count") or 1))
        if str(row.get("verdict") or "") == "correct":
            solved.add(idx)
        parsed = _json_obj(row.get("ai_feedback_json"))
        if parsed:
            feedback[idx] = parsed
    current = next((i for i in range(len(questions)) if i not in solved), None)
    return {
        "attempt_id": int(attempt.get("id") or 0),
        "status": str(attempt.get("status") or "active"),
        "title": attempt.get("title"),
        "source_type": attempt.get("source_type"),
        "source_id": attempt.get("source_id"),
        "total_questions": len(questions),
        "solved_count": len(solved),
        "current_index": current,
        "is_finished": current is None,
        "has_timer": False,
        "questions": [_question_for_student(q) for q in questions],
        "solved_indexes": sorted(solved),
        "tries": tries,
        "last_feedback": feedback.get(current) if current is not None else None,
        "correct_count": int(attempt.get("correct_count") or 0),
        "wrong_count": int(attempt.get("wrong_count") or 0),
        "retry_count": int(attempt.get("retry_count") or 0),
        "dpoints_delta": float(attempt.get("dpoints_delta") or 0.0),
    }


def _questions_from_source(user: dict, source_type: str, source_id: int | None, homework_id: int | None) -> tuple[list[dict], str]:
    source_type = str(source_type or "library_test").strip().lower()
    if source_type == "homework":
        hid = int(homework_id or source_id or 0)
        if hid <= 0:
            raise HTTPException(status_code=400, detail="Homework tanlanmadi")
        homework = _safe(lambda: dbm.get_homework(hid))
        if not homework:
            raise HTTPException(status_code=404, detail="Homework topilmadi")
        test = _safe(lambda: dbm.get_content_test("homework", hid))
        questions = _normalize_questions((test or {}).get("questions"))
        if not questions:
            raise HTTPException(status_code=404, detail="Bu vazifada AI test yo'q")
        return questions, str(homework.get("title") or "Vazifa testi")
    if source_type == "library_test":
        nid = int(source_id or 0)
        node = _node_or_404(nid)
        if str(node.get("kind") or "") != "test":
            raise HTTPException(status_code=400, detail="Bu element test emas")
        if not node.get("is_public"):
            raise HTTPException(status_code=403, detail="Bu test siz uchun ochiq emas")
        questions = _normalize_questions(_json_obj(node.get("payload_json")).get("questions"))
        if not questions:
            raise HTTPException(status_code=404, detail="Testda savol yo'q")
        return questions, str(node.get("title") or "Kutubxona testi")
    if source_type == "weekly_review":
        questions = _weekly_review_questions(int(user.get("id") or 0))
        if not questions:
            raise HTTPException(status_code=404, detail="Haftalik takrorlash uchun material topilmadi")
        return questions, "Haftalik takrorlash"
    raise HTTPException(status_code=400, detail="Yaroqsiz test manbasi")


@router.get("/student/ai-tests/active")
async def ai_test_active(authorization: str | None = Header(default=None)):
    """Ilova qayta ochilganda: active attempt bor-yo'qligini tekshiradi."""
    user = _auth(authorization, STUDENT_ROLES)
    attempt = _safe(lambda: dbm.get_active_ai_test_attempt(int(user.get("id") or 0)))
    if not attempt:
        return {"active": False, "attempt": None}
    return {"active": True, "attempt": _attempt_state(attempt)}


@router.post("/student/ai-tests/start")
async def ai_test_start(payload: AiTestStartRequest, authorization: str | None = Header(default=None)):
    """Testni boshlaydi yoki davom ettiradi.

    Agar shu manba (homework/test) uchun 5 soat ichida boshlangan tugallanmagan
    urinish bo'lsa — o'sha joyidan davom ettiriladi (student chiqib ketib qaytsa
    ma'lumotlari saqlanadi). 5 soatdan oshgan yoki boshqa manba bo'lsa — yangidan."""
    user = _auth(authorization, STUDENT_ROLES)
    from backend.main import _require_student_learning_access

    _require_student_learning_access(user)
    user_id = int(user.get("id") or 0)
    source_id = int(payload.source_id or payload.homework_id or 0) or None

    # 5 soatlik davom ettirish: shu manba uchun active attempt bo'lsa qaytaramiz.
    existing = _safe(lambda: dbm.get_active_ai_test_attempt(user_id))
    if existing and not existing.get("is_finished"):
        same_source = (
            str(existing.get("source_type") or "") == str(payload.source_type)
            and int(existing.get("source_id") or 0) == int(source_id or 0)
        )
        if same_source and _attempt_within_hours(existing.get("started_at"), 5):
            state = _attempt_state(existing)
            if not state["is_finished"]:
                return {"message": "Davom ettirildi", "attempt": state, "resumed": True}

    questions, title = _questions_from_source(user, payload.source_type, payload.source_id, payload.homework_id)
    # Polimorf savollarni (word_practice) shu student uchun random turga, guruh
    # tilidagi (yevro/rus -> ruscha) ko'rsatma bilan ochamiz.
    lang = _student_lang(user)
    study_lang = _study_language_name(_student_subject(user))
    questions = _expand_polymorphic_questions(questions, lang, study_lang)
    # Audio yuklanmagan tinglash mashqlari butun testni bloklamasin — ularni
    # o'tkazib yuboramiz, qolgan mashqlar ishlayveradi.
    questions = [q for q in questions if not q.get("needs_audio_upload")]
    if not questions:
        raise HTTPException(status_code=422, detail="Bu testda ishlaydigan mashq yo'q. O'qituvchiga murojaat qiling.")
    attempt = _safe(
        lambda: dbm.start_ai_test_attempt(
            user_id,
            payload.source_type,
            source_id,
            questions,
            title=title,
        )
    )
    if not attempt:
        raise HTTPException(status_code=500, detail="Test boshlanmadi")
    return {"message": "Test boshlandi", "attempt": _attempt_state(attempt)}


def _attempt_within_hours(started_at: Any, hours: int) -> bool:
    """started_at (timestamp) hozirdan `hours` soat ichida bo'lsa True."""
    if not started_at:
        return False
    from datetime import datetime, timezone

    raw = str(started_at)
    for fmt in ("%Y-%m-%d %H:%M:%S.%f", "%Y-%m-%d %H:%M:%S", "%Y-%m-%dT%H:%M:%S.%f", "%Y-%m-%dT%H:%M:%S"):
        try:
            dt = datetime.strptime(raw[:26] if "." in raw else raw[:19], fmt)
            delta_sec = (datetime.utcnow() - dt).total_seconds()
            # Faqat yuqori chegara: tz farqi delta'ni manfiy qilishi mumkin (baribir
            # yaqinda boshlangan), 5 soatdan oshgan bo'lsa muddati o'tgan.
            return delta_sec <= hours * 3600
        except Exception:
            continue
    return True  # sanani o'qiy olmasak, ehtiyot chorasi sifatida davom ettiramiz


@router.post("/student/ai-tests/{attempt_id}/abandon")
async def ai_test_abandon(attempt_id: int, authorization: str | None = Header(default=None)):
    """Student testdan chiqsa — attempt bekor, keyingi kirishda boshidan."""
    user = _auth(authorization, STUDENT_ROLES)
    attempt = _safe(lambda: dbm.get_ai_test_attempt(int(attempt_id), int(user.get("id") or 0)))
    if not attempt:
        raise HTTPException(status_code=404, detail="Attempt topilmadi")
    dbm.abandon_active_ai_test_attempts(int(user.get("id") or 0))
    return {"message": "Test bekor qilindi, keyingi kirishda boshidan boshlanadi"}


@router.post("/student/ai-tests/upload-audio")
async def ai_test_upload_audio(
    file: UploadFile = File(...),
    authorization: str | None = Header(default=None),
):
    """Speaking javoblari uchun audio yuklash (mikrofon yozuvi)."""
    user = _auth(authorization, STUDENT_ROLES)
    from backend.main import HOMEWORK_UPLOAD_DIR, _upload_web_file

    filename, _ = await _upload_web_file(
        file, HOMEWORK_UPLOAD_DIR, f"aitest_voice_{int(user.get('id') or 0)}", max_size_mb=10
    )
    return {"url": f"/homework/voices/{filename}", "filename": filename}


@router.post("/student/ai-tests/{attempt_id}/answer")
async def ai_test_answer(
    attempt_id: int,
    payload: AiTestAnswerRequest,
    authorization: str | None = Header(default=None),
    x_language: str | None = Header(default=None, alias="X-Language"),
):
    """Bitta savolga javob — darhol tekshiriladi va natija qaytadi."""
    user = _auth(authorization, STUDENT_ROLES)
    user_id = int(user.get("id") or 0)
    attempt = _safe(lambda: dbm.get_ai_test_attempt(int(attempt_id), user_id))
    if not attempt:
        raise HTTPException(status_code=404, detail="Attempt topilmadi")
    if str(attempt.get("status") or "") != "active":
        raise HTTPException(status_code=409, detail="Bu test yakunlangan yoki bekor qilingan")

    state = _attempt_state(attempt)
    index = int(payload.question_index)
    questions = attempt.get("questions") or []
    if index < 0 or index >= len(questions):
        raise HTTPException(status_code=400, detail="Savol raqami noto'g'ri")
    if index in set(state["solved_indexes"]):
        raise HTTPException(status_code=409, detail="Bu savol allaqachon yechilgan")
    if state["current_index"] is not None and index != state["current_index"]:
        raise HTTPException(
            status_code=409,
            detail="Avvalgi savolni tugatmaguncha keyingisiga o'tib bo'lmaydi",
        )

    question = questions[index]
    kind = str(question.get("kind") or "")
    meta = AI_TEST_TYPES.get(kind, {})
    subject = _student_subject(user)

    if meta.get("check") == "auto":
        verdict, feedback = _check_auto(question, payload)
    else:
        verdict, feedback = await _check_with_ai(question, payload, subject, x_language)

    try:
        saved = dbm.save_ai_test_answer(
            int(attempt_id),
            index,
            kind=kind,
            answer_text=payload.answer_text,
            audio_url=payload.audio_url,
            verdict=verdict,
            ai_feedback=feedback,
        )
    except ValueError as exc:
        raise HTTPException(status_code=409, detail=str(exc))

    try_count = int(saved.get("try_count") or 1)
    retry_until_correct = bool(meta.get("retry_until_correct", True))

    dpoints = 0.0
    if verdict == "correct":
        dpoints = _setting("ai_test_correct_reward", 2.0)
        if try_count > 1:
            dpoints -= _setting("ai_test_retry_penalty", 0.5) * (try_count - 1)
        dbm.bump_ai_test_counters(int(attempt_id), correct=1, retries=max(0, try_count - 1), dpoints=dpoints)
    elif not retry_until_correct:
        dpoints = -abs(_setting("ai_test_skip_penalty", 2.0))
        dbm.bump_ai_test_counters(int(attempt_id), wrong=1, retries=max(0, try_count - 1), dpoints=dpoints)
    else:
        # retry_until_correct: cheksiz urinish, jarima yo'q — faqat to'g'risini kutamiz.
        dbm.bump_ai_test_counters(int(attempt_id), retries=1)

    if abs(dpoints) > 0:
        _safe(
            lambda: dbm.add_dpoints(
                user_id, float(dpoints), subject=subject, change_type=f"ai_test_{kind}"
            )
        )

    moved_on = verdict == "correct" or not retry_until_correct

    fresh = _safe(lambda: dbm.get_ai_test_attempt(int(attempt_id), user_id)) or attempt
    new_state = _attempt_state(fresh)
    finished = new_state["current_index"] is None
    if finished:
        _finish_attempt(fresh, user, subject)
        fresh = _safe(lambda: dbm.get_ai_test_attempt(int(attempt_id), user_id)) or fresh
        new_state = _attempt_state(fresh)

    return {
        "verdict": verdict,
        "moved_on": moved_on,
        "try_count": try_count,
        "feedback": feedback,
        "dpoints_delta": round(dpoints, 2),
        "attempt": new_state,
        "finished": finished,
    }


def _student_subject(user: dict) -> str:
    from backend.main import _user_subjects_from_row

    subjects = _safe(lambda: _user_subjects_from_row(user), []) or []
    return str(subjects[0]) if subjects else "English"


def _student_lang(user: dict) -> str:
    """Studentning o'quv tili: guruhi yevro/rus (lang='ru') bo'lsa 'Russian',
    aks holda fandan avtomatik ('Russian'/'Uzbek')."""
    groups = _safe(lambda: dbm.get_user_groups(int(user.get("id") or 0)), []) or []
    for g in groups:
        if str((g or {}).get("lang") or "").strip().lower() == "ru":
            return "Russian"
    return _instruction_language_name(None, _student_subject(user))


def _finish_attempt(attempt: dict, user: dict, subject: str) -> None:
    dbm.bump_ai_test_counters(int(attempt.get("id") or 0), complete=True)
    source_type = str(attempt.get("source_type") or "")
    if source_type == "homework":
        hid = int(attempt.get("source_id") or 0)
        if hid > 0:
            _safe(
                lambda: dbm.upsert_homework_submission(
                    homework_id=hid,
                    student_id=int(user.get("id") or 0),
                    status="pending_review",
                    note="AI test yakunlandi",
                    proof_image_url=None,
                    proof_images_json=None,
                )
            )
    if source_type == "weekly_review":
        week_start = _week_bounds()[0]
        _safe(lambda: dbm.complete_weekly_review(int(user.get("id") or 0), week_start))
        reward = _setting("weekly_review_reward", 10.0)
        if reward:
            _safe(
                lambda: dbm.add_dpoints(
                    int(user.get("id") or 0), float(reward), subject=subject, change_type="weekly_review_done"
                )
            )


# ── Avtomatik tekshirish ────────────────────────────────────────────────────

def _norm_text(value: Any) -> str:
    text = str(value or "").strip().lower()
    text = re.sub(r"[^\w\s'’-]", " ", text, flags=re.UNICODE)
    return re.sub(r"\s+", " ", text).strip()


def _check_auto(question: dict, payload: AiTestAnswerRequest) -> tuple[str, dict]:
    kind = str(question.get("kind") or "")
    if kind == "listening":
        options = question.get("options") or []
        correct = int(question.get("correct_index") or 0)
        chosen = payload.choice_index
        if chosen is None:
            return "wrong", {"reason": "Javob tanlanmadi"}
        ok = int(chosen) == correct
        return ("correct" if ok else "wrong"), {
            "correct_answer": options[correct] if 0 <= correct < len(options) else None,
        }
    if kind == "matching":
        pairs = {str(p.get("left")): str(p.get("right")) for p in (question.get("pairs") or [])}
        given = {str(k): str(v) for k, v in (payload.pairs or {}).items()}
        wrong_keys = [left for left, right in pairs.items() if _norm_text(given.get(left)) != _norm_text(right)]
        return ("correct" if not wrong_keys else "wrong"), {
            "wrong_items": wrong_keys,
            "wrong_count": len(wrong_keys),
        }
    if kind == "scrambled_sentence":
        expected = _norm_text(question.get("answer"))
        given = _norm_text(" ".join(payload.order or []) or payload.answer_text)
        return ("correct" if given == expected else "wrong"), {
            "hint": "So'zlar tartibini qayta ko'ring" if given != expected else None,
        }
    if kind == "passage_cloze":
        blanks = question.get("blanks") or []
        given_list = payload.blanks or []
        wrong_positions = []
        for i, b in enumerate(blanks):
            expected_set = {_norm_text(b.get("answer")), *[_norm_text(x) for x in (b.get("accepted_answers") or [])]}
            expected_set = {e for e in expected_set if e}
            given = _norm_text(given_list[i]) if i < len(given_list) else ""
            if given not in expected_set:
                wrong_positions.append(i + 1)
        if not wrong_positions:
            return "correct", {}
        return "wrong", {
            "wrong_positions": wrong_positions,
            "wrong_count": len(wrong_positions),
            "hint": f"{len(wrong_positions)} ta bo'sh joy noto'g'ri",
        }
    if kind == "reading_set":
        subs = question.get("sub_questions") or []
        given_list = payload.blanks or []
        wrong_positions = []
        for i, s in enumerate(subs):
            expected_set = {_norm_text(s.get("answer")), *[_norm_text(x) for x in (s.get("accepted_answers") or [])]}
            expected_set = {e for e in expected_set if e}
            given = _norm_text(given_list[i]) if i < len(given_list) else ""
            if given not in expected_set:
                wrong_positions.append(i + 1)
        if not wrong_positions:
            return "correct", {}
        return "wrong", {
            "wrong_positions": wrong_positions,
            "wrong_count": len(wrong_positions),
            "hint": f"{len(wrong_positions)} ta savol noto'g'ri",
        }
    expected = [question.get("answer"), *(question.get("accepted_answers") or [])]
    normalized = {_norm_text(e) for e in expected if e}
    given = _norm_text(payload.answer_text)
    if not given:
        return "wrong", {"reason": "Javob bo'sh"}
    if given in normalized:
        return "correct", {}
    close = any(_levenshtein(given, e) <= 1 for e in normalized)
    return "wrong", {"almost": close, "hint": "Imloni tekshiring" if close else None}


def _levenshtein(a: str, b: str) -> int:
    if a == b:
        return 0
    if abs(len(a) - len(b)) > 2:
        return 99
    prev = list(range(len(b) + 1))
    for i, ca in enumerate(a, 1):
        cur = [i]
        for j, cb in enumerate(b, 1):
            cur.append(min(prev[j] + 1, cur[j - 1] + 1, prev[j - 1] + (ca != cb)))
        prev = cur
    return prev[-1]


# ── AI tekshirish (grammatika + talaffuz + daraja mosligi) ──────────────────

async def _transcribe(audio_url: str, subject: str, x_language: str | None) -> str:
    """xAI STT bilan audio -> matn."""
    from backend.main import HOMEWORK_UPLOAD_DIR, _speech_language_code
    from ai_generator import _get_xai_api_key
    import aiohttp

    api_key = _safe(lambda: _get_xai_api_key())
    if not api_key:
        raise HTTPException(status_code=503, detail="AI servis sozlanmagan")
    filename = Path(urlparse(str(audio_url or "")).path).name
    audio_path = HOMEWORK_UPLOAD_DIR / filename
    if not filename or not audio_path.is_file():
        raise HTTPException(status_code=404, detail="Audio fayl serverda topilmadi")
    form = aiohttp.FormData()
    form.add_field("language", _speech_language_code(subject, x_language))
    form.add_field("format", "true")
    content_type = mimetypes.guess_type(audio_path.name)[0] or "application/octet-stream"
    with audio_path.open("rb") as fh:
        # Fayl maydoni xAI STT uchun oxirgi bo'lishi shart.
        form.add_field("file", fh, filename=audio_path.name, content_type=content_type)
        async with aiohttp.ClientSession() as session:
            async with session.post(
                "https://api.x.ai/v1/stt",
                headers={"Authorization": f"Bearer {api_key}"},
                data=form,
                timeout=aiohttp.ClientTimeout(total=90),
            ) as resp:
                if resp.status != 200:
                    detail = (await resp.text())[:300]
                    logger.warning("ai-test stt failed status=%s detail=%s", resp.status, detail)
                    raise HTTPException(status_code=503, detail="Audio transkript qilinmadi, qayta urinib ko'ring")
                data = await resp.json(content_type=None)
                return str((data or {}).get("text") or "").strip()


def _ai_check_prompt(question: dict, answer: str, subject: str, result_language: str, spoken: bool) -> str:
    kind = str(question.get("kind") or "")
    level = str(question.get("target_level") or question.get("level") or "").strip()
    parts = [
        f"You are a strict but fair {subject} teacher checking one student answer.",
        f"Exercise type: {kind}.",
        f"Task given to the student: {question.get('prompt') or ''}",
    ]
    if question.get("word"):
        parts.append(f"The student MUST use this word: {question['word']}")
    if question.get("passage"):
        parts.append(f"Reference text:\n{question['passage']}")
    if question.get("reference_answer"):
        parts.append(f"A model answer (not the only valid one): {question['reference_answer']}")
    if level:
        parts.append(
            f"The student's level is {level}. The answer must be close to that level: "
            "not a trivial 2-word sentence for B1+, and not over-complicated for A1/A2."
        )
    parts.append(
        ("This answer was SPOKEN and transcribed by speech-to-text. "
         "Judge pronunciation from the transcript: if words came out as different words, "
         "flag them as pronunciation problems. Ignore missing punctuation and capitalisation."
         if spoken else
         "This answer was TYPED. Judge spelling, punctuation and grammar.")
    )
    parts.append(
        "Mark the answer as correct ONLY if: grammar is correct, the required word is used properly "
        "(if one was required), the meaning answers the task, and the level fits."
    )
    parts.append(
        f"Write every explanation field in {result_language}. Do not translate the student's own answer.\n"
        "Return ONLY valid JSON, no markdown:\n"
        '{"is_correct":true,'
        '"grammar_ok":true,"pronunciation_ok":true,"level_ok":true,'
        '"corrected":"the corrected version of the answer",'
        '"grammar_errors":[{"original":"...","correction":"...","explanation":"..."}],'
        '"pronunciation_errors":[{"word":"...","note":"..."}],'
        '"score":8.5,"feedback":"one short encouraging sentence with the main fix"}'
    )
    parts.append(f"Student answer:\n{answer}")
    return "\n\n".join(parts)


async def _check_with_ai(
    question: dict, payload: AiTestAnswerRequest, subject: str, x_language: str | None
) -> tuple[str, dict]:
    from backend.main import _ai_result_language

    kind = str(question.get("kind") or "")
    meta = AI_TEST_TYPES.get(kind, {})
    input_mode = str(meta.get("input") or "text")
    spoken = False
    answer = str(payload.answer_text or "").strip()

    if input_mode == "audio" or (input_mode == "audio_or_text" and payload.audio_url):
        if not payload.audio_url:
            raise HTTPException(status_code=400, detail="Ovozli javob yuborilmadi")
        answer = await _transcribe(payload.audio_url, subject, x_language)
        spoken = True
        if not answer:
            return "wrong", {
                "transcript": "",
                "feedback": "Ovozingiz aniq eshitilmadi. Yana bir marta, sekinroq va balandroq gapiring.",
                "pronunciation_ok": False,
            }
    if not answer:
        return "wrong", {"feedback": "Javob bo'sh qoldirilgan."}

    from ai_generator import XAI_ENDPOINT, _get_xai_api_key, _xai_apply_payload_tuning, get_grok_model_candidates
    import aiohttp

    api_key = _safe(lambda: _get_xai_api_key())
    if not api_key:
        raise HTTPException(status_code=503, detail="AI servis sozlanmagan")
    prompt = _ai_check_prompt(question, answer, subject, _ai_result_language(x_language), spoken)
    raw = ""
    try:
        async with aiohttp.ClientSession() as session:
            for model in get_grok_model_candidates()[:2]:
                body = {
                    "model": model,
                    "messages": [
                        {"role": "system", "content": "You are a precise language-exercise grader. Output only JSON."},
                        {"role": "user", "content": prompt},
                    ],
                    "temperature": 0.2,
                    "max_tokens": 1200,
                }
                _xai_apply_payload_tuning(body, model=model, stream=False)
                async with session.post(
                    XAI_ENDPOINT,
                    headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
                    json=body,
                    timeout=aiohttp.ClientTimeout(total=60),
                ) as resp:
                    if resp.status == 200:
                        data = await resp.json(content_type=None)
                        raw = str(data.get("choices", [{}])[0].get("message", {}).get("content", "")).strip()
                        if raw:
                            break
    except Exception as exc:
        logger.exception("ai-test grading failed kind=%s: %s", kind, exc)
        raise HTTPException(status_code=503, detail="AI tekshirish vaqtincha ishlamayapti, qayta urinib ko'ring")

    parsed = _extract_json_object(raw)
    if not parsed:
        raise HTTPException(status_code=503, detail="AI javobi tushunarsiz, qayta urinib ko'ring")
    verdict = "correct" if bool(parsed.get("is_correct")) else "wrong"
    parsed["transcript"] = answer if spoken else None
    parsed["was_spoken"] = spoken
    return verdict, parsed


# ═══════════════════════════════════════════════════════════════════════════
# 4. HAFTALIK MAJBURIY TAKRORIY TEST
# ═══════════════════════════════════════════════════════════════════════════

def _week_bounds(now: datetime | None = None) -> tuple[str, str]:
    """Joriy haftaning dushanba–yakshanba chegaralari (YYYY-MM-DD)."""
    current = now or datetime.now()
    monday = current - timedelta(days=current.weekday())
    sunday = monday + timedelta(days=6)
    return monday.strftime("%Y-%m-%d"), sunday.strftime("%Y-%m-%d")


def _weekly_review_questions(user_id: int, limit: int = 0) -> list[dict]:
    """Hafta davomida o'quvchiga berilgan BARCHA mashqlarni yaxlit yig'adi.

    Manbalar: hafta ichida ishlangan content/AI testlar + shu hafta berilgan
    homeworklarga biriktirilgan testlar (kutubxonadan berilganlar ham).
    `limit=0` — cheklov yo'q (hammasi bir homeworkda beriladi).
    """
    week_start, week_end = _week_bounds()
    collected: list[dict] = []
    seen: set[str] = set()

    def _add_all(source: list[dict]) -> None:
        for question in source or []:
            key = f"{question.get('kind')}|{_norm_text(question.get('prompt') or question.get('word') or question.get('passage'))}"
            if not key.strip("|") or key in seen:
                continue
            if question.get("needs_audio_upload"):
                continue
            seen.add(key)
            collected.append(question)

    # 1) Hafta ichida ishlangan urinishlar (content_test + ai_test).
    rows = _safe(
        lambda: dbm.list_user_week_attempts(user_id, f"{week_start} 00:00:00", f"{week_end} 23:59:59"), []
    ) or []
    for row in rows:
        if str(row.get("origin") or "") == "ai_test":
            attempt = _safe(lambda aid=int(row.get("id") or 0): dbm.get_ai_test_attempt(aid, user_id))
            _add_all((attempt or {}).get("questions") or [])
        else:
            test = _safe(
                lambda: dbm.get_content_test(
                    str(row.get("content_type") or "book"), int(row.get("content_id") or 0)
                )
            )
            _add_all(_normalize_questions((test or {}).get("questions")))

    # 2) Shu hafta berilgan homeworklardagi testlar (ishlanmagan bo'lsa ham).
    for hw in _student_week_homeworks(user_id, week_start, week_end):
        hid = int(hw.get("id") or 0)
        if hid <= 0:
            continue
        test = _safe(lambda h=hid: dbm.get_content_test("homework", h))
        _add_all(_normalize_questions((test or {}).get("questions")))

    return collected if limit <= 0 else collected[:limit]


def _student_week_homeworks(user_id: int, week_start: str, week_end: str) -> list[dict]:
    """Shu hafta o'quvchiga berilgan homeworklar (shaxsiy + guruh)."""
    conn = dbm.get_conn()
    cur = conn.cursor()
    try:
        cur.execute(
            """
            SELECT h.id
            FROM web_homeworks h
            WHERE COALESCE(h.status,'active') <> 'deleted'
              AND h.created_at >= ? AND h.created_at <= ?
              AND (
                h.student_id = ?
                OR (h.group_id IS NOT NULL AND h.group_id IN (
                     SELECT group_id FROM user_groups WHERE user_id = ?
                   ))
              )
            ORDER BY h.id DESC
            LIMIT 200
            """,
            (f"{week_start} 00:00:00", f"{week_end} 23:59:59", int(user_id), int(user_id)),
        )
        return [dbm._row_to_dict(r) for r in (cur.fetchall() or [])]
    except Exception:
        logger.exception("weekly review: homework list failed user_id=%s", user_id)
        return []
    finally:
        conn.close()


@router.get("/student/weekly-review")
async def student_weekly_review(authorization: str | None = Header(default=None)):
    """Studentga haftalik takrorlash holati."""
    user = _auth(authorization, STUDENT_ROLES)
    week_start, week_end = _week_bounds()
    review = _safe(lambda: dbm.get_weekly_review(int(user.get("id") or 0), week_start))
    available = len(_weekly_review_questions(int(user.get("id") or 0)))
    return {
        "week_start": week_start,
        "week_end": week_end,
        "status": str((review or {}).get("status") or "none"),
        "homework_id": (review or {}).get("homework_id"),
        "question_count": available,
        "is_required": bool(review) and str((review or {}).get("status") or "") == "assigned",
    }


@router.post("/admin/weekly-review/run")
async def admin_weekly_review_run(
    authorization: str | None = Header(default=None),
    dry_run: bool = Query(default=False),
):
    """Haftalik takroriy testni barcha studentlarga homework qilib tarqatadi.

    Cron / scheduler shu endpointni hafta oxirida chaqiradi."""
    _auth(authorization, {"admin", "superadmin"})
    result = await asyncio.to_thread(_run_weekly_review_job, bool(dry_run))
    return result


def _run_weekly_review_job(dry_run: bool = False) -> dict:
    week_start, week_end = _week_bounds()
    conn = dbm.get_conn()
    cur = conn.cursor()
    try:
        cur.execute(
            """
            SELECT DISTINCT u.id
            FROM users u
            JOIN user_groups ug ON ug.user_id = u.id
            WHERE u.login_type IN (1, 2)
            """
        )
        student_ids = [int(r["id"]) for r in (cur.fetchall() or []) if r and r.get("id")]
    except Exception:
        logger.exception("weekly review: student list failed")
        student_ids = []
    finally:
        conn.close()

    assigned = 0
    skipped = 0
    for student_id in student_ids:
        existing = _safe(lambda sid=student_id: dbm.get_weekly_review(sid, week_start))
        if existing:
            skipped += 1
            continue
        questions = _weekly_review_questions(student_id)
        if not questions:
            skipped += 1
            continue
        if dry_run:
            assigned += 1
            continue
        teacher_id = _weekly_review_teacher_id(student_id)
        # Deadline: berilgan paytdan 24 soat.
        due_at = (datetime.now() + timedelta(hours=24)).strftime("%Y-%m-%d %H:%M:00")
        homework = _safe(
            lambda sid=student_id, tid=teacher_id: dbm.create_homework(
                tid,
                sid,
                f"Haftalik takrorlash ({week_start} — {week_end})",
                description=(
                    "Bu hafta berilgan barcha mashqlar bitta vazifada takrorlanadi. "
                    "Bajarish majburiy — analitika foizingizga ta'sir qiladi. "
                    "Muddat: 24 soat."
                ),
                due_at=due_at,
                homework_kind="test",
            )
        )
        homework_id = int((homework or {}).get("id") or 0)
        if homework_id > 0:
            _safe(
                lambda hid=homework_id, tid=teacher_id: dbm.save_content_test(
                    "homework", hid, json.dumps(questions, ensure_ascii=False), tid,
                    title="Haftalik takrorlash", raw_questions=True,
                )
            )
        _safe(
            lambda sid=student_id, hid=homework_id: dbm.upsert_weekly_review(
                sid, week_start, week_end, hid or None, attempt_count=len(questions)
            )
        )
        # O'quvchiga bildirishnoma (web + push).
        if homework_id > 0:
            try:
                from backend.main import (
                    _notify_homework_created,
                    _homework_audience_student_ids,
                    _attach_homework_runtime_fields,
                )

                hw = _safe(lambda h=homework_id: dbm.get_homework(h))
                if hw:
                    _safe(lambda: _attach_homework_runtime_fields(hw))
                    teacher = _safe(lambda t=teacher_id: dbm.get_user_by_id(int(t))) or {}
                    import asyncio as _asyncio

                    _asyncio.run(
                        _notify_homework_created(hw, teacher, _homework_audience_student_ids(hw))
                    )
            except Exception:
                logger.exception("weekly review notify failed homework_id=%s", homework_id)
        assigned += 1
    logger.info("weekly review job week=%s assigned=%s skipped=%s dry_run=%s", week_start, assigned, skipped, dry_run)
    return {
        "week_start": week_start,
        "week_end": week_end,
        "students": len(student_ids),
        "assigned": assigned,
        "skipped": skipped,
        "dry_run": dry_run,
    }


def _weekly_review_teacher_id(student_id: int) -> int:
    """Studentning guruh o'qituvchisi; topilmasa 0 (tizim vazifasi)."""
    conn = dbm.get_conn()
    cur = conn.cursor()
    try:
        cur.execute(
            """
            SELECT g.teacher_id
            FROM user_groups ug
            JOIN groups g ON g.id = ug.group_id
            WHERE ug.user_id = ? AND g.teacher_id IS NOT NULL
            ORDER BY ug.group_id DESC LIMIT 1
            """,
            (int(student_id),),
        )
        row = cur.fetchone()
        return int((row or {}).get("teacher_id") or 0)
    except Exception:
        return 0
    finally:
        conn.close()


@router.post("/admin/weekly-review/penalize")
async def admin_weekly_review_penalize(authorization: str | None = Header(default=None)):
    """Muddati (24 soat) o'tgan takrorlashni bajarmaganlarga jarima (D'point)."""
    _auth(authorization, {"admin", "superadmin"})
    now = datetime.now()
    week_start = (now - timedelta(days=now.weekday())).strftime("%Y-%m-%d")
    penalty = abs(_setting("weekly_review_missed_penalty", 5.0))
    dbm.ensure_weekly_reviews_schema()
    conn = dbm.get_conn()
    cur = conn.cursor()
    try:
        cur.execute(
            "SELECT user_id FROM weekly_reviews WHERE week_start=? AND status='assigned'",
            (week_start,),
        )
        user_ids = [int(r["user_id"]) for r in (cur.fetchall() or []) if r and r.get("user_id")]
    finally:
        conn.close()
    if penalty > 0:
        for uid in user_ids:
            _safe(lambda u=uid: dbm.add_dpoints(u, -penalty, change_type="weekly_review_missed"))
    return {"week_start": week_start, "penalized": len(user_ids), "penalty": penalty}
